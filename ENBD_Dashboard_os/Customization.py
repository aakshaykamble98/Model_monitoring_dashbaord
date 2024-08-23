import streamlit as st
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import base64
import os

# Function to create PowerPoint presentation
def create_ppt(bg_color, font_color, ribbon_color_1, ribbon_color_2, row_bg_color, row_font_color, content_font_color, title_font_color, ribbon_font_1, ribbon_font_2, ribbon_font_color):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # Using a blank slide layout
    
    # Slide with background and font color
    slide1 = prs.slides.add_slide(slide_layout)
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(bg_color[1:])
    
    title = slide1.shapes.title
    title.text = "Sample Title"
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor.from_string(font_color[1:])
            
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    ribbon_height = Inches(0.28)
    half_ribbon_width = slide_width / 2
    ribbon_top = slide_height - ribbon_height
            
    slide_number_box = slide1.shapes.add_textbox(half_ribbon_width, ribbon_top, half_ribbon_width, ribbon_height)
    text_frame = slide_number_box.text_frame
    text_frame.text = "1"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor.from_string(ribbon_font_color[1:])
    p.alignment = PP_ALIGN.RIGHT
    
    # Slide with ribbons and table
    slide2 = prs.slides.add_slide(slide_layout)
    
    title = slide2.shapes.title
    title.text = "Sample Title"
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor.from_string(title_font_color[1:])
            
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # First half of the ribbon
    shape_amber = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, ribbon_top, half_ribbon_width, ribbon_height
    )
    fill_amber = shape_amber.fill
    fill_amber.solid()
    fill_amber.fore_color.rgb = RGBColor.from_string(ribbon_color_1[1:])
    
    line_amber = shape_amber.line
    line_amber.color.rgb = RGBColor.from_string(ribbon_color_1[1:])
    
    text_frame_amber = shape_amber.text_frame
    text_frame_amber.text = "\tModel Monitoring"
    text_frame_amber.paragraphs[0].font.size = Pt(9)
    text_frame_amber.paragraphs[0].font.color.rgb = RGBColor.from_string(ribbon_font_1[1:])
    
    # Second half of the ribbon
    shape_blue = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, half_ribbon_width, ribbon_top, half_ribbon_width, ribbon_height
    )
    fill_blue = shape_blue.fill
    fill_blue.solid()
    fill_blue.fore_color.rgb = RGBColor.from_string(ribbon_color_2[1:])
    
    line_blue = shape_blue.line
    line_blue.color.rgb = RGBColor.from_string(ribbon_color_2[1:])
    
    slide_number_box = slide2.shapes.add_textbox(half_ribbon_width, ribbon_top, half_ribbon_width, ribbon_height)
    text_frame = slide_number_box.text_frame
    text_frame.text = "2"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor.from_string(ribbon_font_2[1:])
    p.alignment = PP_ALIGN.RIGHT
    
    # Add logo image to the upper right corner
    base_dir = os.path.dirname(__file__)
    
    # Logo path
    logo_path = os.path.join(base_dir, 'Images', 'ENBD_s.jpg')  # Replace with your actual file path
    logo_left = slide_width - Inches(0.6)
    logo_top = Inches(0.15)
    logo_height = Inches(0.45)
    slide2.shapes.add_picture(logo_path, logo_left, logo_top, height=logo_height)
    
    # Add smaller table
    rows, cols = 3, 3
    left = Inches(2.7)
    top = Inches(2)
    width = Inches(4.7)  # Adjusted width
    height = Inches(3)  # Adjusted height
    table = slide2.shapes.add_table(rows, cols, left, top, width, height).table
    table_style_id = table._tbl.tblPr.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}tableStyleId"
    )
    table_style_id.text = "{5940675A-B579-460E-94D1-54222C63F5DA}"
    
    # Set font size for the table cells
    font_size = Pt(15)  # Smaller font size
    
    # Set table headers
    headers = ["Header 1", "Header 2", "Header 3"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(row_bg_color[1:])
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = RGBColor.from_string(row_font_color[1:])
            paragraph.font.size = font_size  # Set font size for headers
    
    # Set table content
    for row in range(1, rows):
        for col in range(cols):
            cell = table.cell(row, col)
            cell.text = f"Data {row+1},{col+1}"
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.color.rgb = RGBColor.from_string(content_font_color[1:])
                paragraph.font.size = font_size  # Set font size for content
                # Adjust padding to fit the smaller cell size
                cell.text_frame.margin_left = Inches(0.5)
                cell.text_frame.margin_right = Inches(0.5)
                cell.text_frame.margin_top = Inches(0.5)
                cell.text_frame.margin_bottom = Inches(0.5)
    
    return prs

# Function to save PowerPoint presentation
def save_ppt(prs):
    ppt_io = BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# Streamlit app
def app():
    st.markdown(
        """
        <div style='background-color: #008080; border-radius: 5px; padding: 1px;'>
            <h1 style='text-align: center; font-size: 28px; color: white;'>PowerPoint Customization</h1>
        </div>

        """
        ,
        unsafe_allow_html=True
    )

    #st.title("PowerPoint Customizer")
    
    # Initialize session state variables if not already set for slide 1
    if 'bg_color' not in st.session_state:
        st.session_state.bg_color = "#008080"
    if 'font_color' not in st.session_state:
        st.session_state.font_color = "#FFFFFF"
    if 'ribbon_font_color' not in st.session_state:
        st.session_state.ribbon_font_color = "#FFFFFF"
        
    # Initialize session state variables if not already set for slide 2 
    if 'title_font_color' not in st.session_state:
        st.session_state.title_font_color = "#000000"
    if 'ribbon_color_1' not in st.session_state:
        st.session_state.ribbon_color_1 = "#FFBF00"
    if 'ribbon_font_1' not in st.session_state:
        st.session_state.ribbon_font_1 = "#000000"
    if 'ribbon_color_2' not in st.session_state:
        st.session_state.ribbon_color_2 = "#FFBF00"
    if 'ribbon_font_2' not in st.session_state:
        st.session_state.ribbon_font_2 = "#FFFFFF"
    if 'row_bg_color' not in st.session_state:
        st.session_state.row_bg_color = "#008080"
    if 'row_font_color' not in st.session_state:
        st.session_state.row_font_color = "#FFFFFF"
    if 'content_font_color' not in st.session_state:
        st.session_state.content_font_color = "#000000"
    
    # Define tabs
    tab1, tab2 = st.tabs(["Title Slide", "Core Slide"])
    
    # Slide 1 Customizer
    with tab1:
        
        # Display the centered heading with space below it
        st.markdown(
                    """
                    <div style='text-align: center;
                    font-size: 22px;
                    margin-bottom: 20px;'>  <!-- Adjust margin-bottom to increase/decrease spacing -->
                        <strong>Background Color and Font Customization</strong>
                    </div>
                    """,
                    unsafe_allow_html=True)
            
        col1, col2, col3 = st.columns([1, 1, 2])
        
        with col1:
            st.session_state.bg_color = st.color_picker("Background color", st.session_state.bg_color, key="bg_color_picker")
            st.session_state.ribbon_font_color = st.color_picker("Ribbon font color", st.session_state.ribbon_font_color, key="ribbon_font_color_picker")

        with col2:
            st.session_state.font_color = st.color_picker("Font color", st.session_state.font_color, key="font_color_picker")
        
        with col3:
            # Live preview of the slide
            st.markdown(f"""
                <div style="
                    width:350px;
                    height:300px;
                    border:1px solid #000;
                    position:relative;
                    background-color:{st.session_state.bg_color};
                    margin-bottom:10px;
                ">
                    <h2 style="
                        color:{st.session_state.font_color};
                        position:absolute;
                        top:50%;
                        left:50%;
                        transform:translate(-50%, -50%);
                        margin:0;
                        text-align:center;
                    ">
                        Sample Title
                    </h2>
                    <div style="
                        width:100%;
                        height:22px;
                        position:absolute;
                        bottom:0;
                        right:0;
                        text-align:center;
                        color:{st.session_state.ribbon_font_color};
                        line-height:28px;
                        font-size:12px;
                        box-sizing:border-box;
                        border-left:1px solid #000;
                    ">
                        1
                    </div>
                    <img src="https://via.placeholder.com/90x65" style="
                        position:absolute;
                        top:10px;
                        right:10px;
                        height:45px;
                    ">
                </div>
            """, unsafe_allow_html=True)
    
    # Slide 2 Customizer
    with tab2:
        # Display the centered heading with space below it
        st.markdown(
                    """
                    <div style='text-align: center;
                    font-size: 22px;
                    margin-bottom: 20px;'>  <!-- Adjust margin-bottom to increase/decrease spacing -->
                        <strong>Ribbons, Title and Table Customization</strong>
                    </div>
                    """,
                    unsafe_allow_html=True)
        
        col1, col2, col3 = st.columns([1, 1, 2])
        
        with col1:
            st.session_state.title_font_color = st.color_picker("Title font color", st.session_state.title_font_color, key="title_font_color_picker")
            st.session_state.row_font_color = st.color_picker("Table first row font color", st.session_state.row_font_color, key="row_font_color_picker")
            st.session_state.ribbon_color_1 = st.color_picker("First ribbon color", st.session_state.ribbon_color_1, key="ribbon_color_1_picker")
            st.session_state.ribbon_color_2 = st.color_picker("Second ribbon color", st.session_state.ribbon_color_2, key="ribbon_color_2_picker")
 
        with col2:
            st.session_state.row_bg_color = st.color_picker("Table first row background color", st.session_state.row_bg_color, key="row_bg_color_picker")
            st.session_state.content_font_color = st.color_picker("Table content font color", st.session_state.content_font_color, key="content_font_color_picker")
            st.session_state.ribbon_font_1 = st.color_picker("First ribbon font color", st.session_state.ribbon_font_1, key="ribbon_font_1_picker")
            st.session_state.ribbon_font_2 = st.color_picker("Second ribbon font color", st.session_state.ribbon_font_2, key="ribbon_font_2_picker")
        
        with col3:
            # Display a live preview of the slide with the selected ribbon and table colors
            st.markdown(f"""
                <div style="
                    width:350px;
                    height:300px;
                    border:1px solid #000;
                    position:relative;
                    background-color:white;
                    overflow:auto;
                ">
                    <h2 style="
                        color:{st.session_state.title_font_color};
                        position:absolute;
                        top:10px;
                        left:50%;
                        transform:translateX(-50%);
                        margin:0;
                        font-size:16px;
                        text-align:center;
                    ">
                        Sample Title
                    </h2>
                    <div style="
                        background-color:{st.session_state.ribbon_color_1};
                        width:50%;
                        height:22px;
                        position:absolute;
                        bottom:0;
                        left:0;
                        text-align:center;
                        color:{st.session_state.ribbon_font_1};
                        line-height:22px;
                        font-size:12px;
                        box-sizing:border-box;
                        border-right:1px solid #000;
                    ">
                        Model Monitoring
                    </div>
                    <div style="
                        background-color:{st.session_state.ribbon_color_2};
                        width:50%;
                        height:22px;
                        position:absolute;
                        bottom:0;
                        right:0;
                        text-align:center;
                        color:{st.session_state.ribbon_font_2};
                        line-height:22px;
                        font-size:7px;
                        box-sizing:border-box;
                        border-left:1px solid #000;
                    ">
                        2
                    </div>
                    <div style="
                        position:absolute;
                        top:65px;
                        left:65px;
                        width:200px;
                        height:135px;  /* Further adjusted height */
                        border:1px solid #000;
                        border-collapse:collapse;
                        font-size:10px; /* Smaller font size */
                        overflow:auto;
                    ">
                        <table style="
                            width:100%;
                            height:100%;
                            border-collapse:collapse;
                            font-size:10px;
                        ">
                            <tr>
                                <th style="
                                    background-color:{st.session_state.row_bg_color};
                                    color:{st.session_state.row_font_color};
                                    padding:2px;
                                    border:1px solid black;
                                    font-size:10px;
                                ">Header 1</th>
                                <th style="
                                    background-color:{st.session_state.row_bg_color};
                                    color:{st.session_state.row_font_color};
                                    padding:2px;
                                    border:1px solid black;
                                    font-size:10px;
                                ">Header 2</th>
                                <th style="
                                    background-color:{st.session_state.row_bg_color};
                                    color:{st.session_state.row_font_color};
                                    padding:2px;
                                    border:1px solid black;
                                    font-size:10px;
                                ">Header 3</th>
                            </tr>
                            <tr>
                                <td style="
                                    color:{st.session_state.content_font_color};
                                    padding:2px;
                                    border:1px solid black;
                                    font-size:10px;
                                ">Data 1,1</td>
                                <td style="
                                    color:{st.session_state.content_font_color};
                                    padding:2px;
                                    border:1px solid black;
                                    font-size:10px;
                                ">Data 1,2</td>
                                <td style="
                                    color:{st.session_state.content_font_color};
                                    padding:2px;
                                    border:1px solid black;
                                    font-size:10px;
                                ">Data 1,3</td>
                            </tr>
                            <tr>
                                <td style="
                                    color:{st.session_state.content_font_color};
                                    padding:2px;
                                    border:1px solid black;
                                    font-size:10px;
                                ">Data 2,1</td>
                                <td style="
                                    color:{st.session_state.content_font_color};
                                    padding:2px;
                                    border:1px solid black;
                                    font-size:10px;
                                ">Data 2,2</td>
                                <td style="
                                    color:{st.session_state.content_font_color};
                                    padding:2px;
                                    border:1px solid black;
                                    font-size:10px;
                                ">Data 2,3</td>
                            </tr>
                            <tr>
                                <td style="
                                    color:{st.session_state.content_font_color};
                                    padding:2px;
                                    border:1px solid black;
                                    font-size:10px;
                                ">Data 3,1</td>
                                <td style="
                                    color:{st.session_state.content_font_color};
                                    padding:2px;
                                    border:1px solid black;
                                    font-size:10px;
                                ">Data 3,2</td>
                                <td style="
                                    color:{st.session_state.content_font_color};
                                    padding:2px;
                                    border:1px solid black;
                                    font-size:10px;
                                ">Data 3,3</td>
                            </tr>
                        </table>
                    </div>
                    <img src="https://via.placeholder.com/90x65" style="
                        position:absolute;
                        top:10px;
                        right:10px;
                        height:45px;
                    ">
                </div>
            """, unsafe_allow_html=True)
        
            prs = create_ppt(st.session_state.bg_color, st.session_state.font_color, st.session_state.ribbon_color_1, st.session_state.ribbon_color_2, st.session_state.row_bg_color, st.session_state.row_font_color, st.session_state.content_font_color, st.session_state.title_font_color, st.session_state.ribbon_font_1, st.session_state.ribbon_font_2, st.session_state.ribbon_font_color)
            ppt_io = save_ppt(prs)
            ppt_bytes = ppt_io.seek(0)
            ppt_bytes = ppt_io.getvalue()
            
            # Base directory of the current script
            base_dir = os.path.dirname(__file__)
            
            # Construct the path for the image
            image_path = os.path.join(base_dir, "Images", "ppt_logo.png")
            
            # Read and encode the image to base64
            with open(image_path, "rb") as image_file:
                image_base64 = base64.b64encode(image_file.read()).decode()
            
            # Encode the PPT data to Base64
            ppt_data_base64 = base64.b64encode(ppt_bytes).decode()
            
            # Help text
            help_text = "Click here to download the Customized PowerPoint Presentation"
            
            custom_css = """
            <style>
                .ppt-download-button {
                    position: absolute;
                    top: 50px;
                    left: -450px;
                    cursor: pointer
                }
            </style>
            """
            st.markdown(custom_css, unsafe_allow_html=True)
            
            # Create the HTML for the button with image and help text
            button_html = f"""
            <a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{ppt_data_base64}" download="Overview.pptx" class="ppt-download-button" title="{help_text}">
                <img src="data:image/png;base64,{image_base64}" alt="Download PPT">
            </a>
            """
            
            # Display the custom download button in Streamlit
            st.markdown(button_html, unsafe_allow_html=True)

if __name__ == "__main__":
    app()