import streamlit as st
from PIL import Image
import io
import pandas as pd
import os
import pickle
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl.formatting.rule import Rule
from openpyxl.styles.differential import DifferentialStyle
import base64
from openpyxl.formatting.rule import FormulaRule, IconSetRule, IconSet, FormatObject
from openpyxl import load_workbook
from openpyxl.styles import PatternFill
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
from pptx.util import Inches, Pt
from pptx import Presentation
import sys
from streamlit_option_menu import option_menu

# # Set the working directory to the root of the project
# project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
# os.chdir(project_root)
# sys.path.append(project_root)

# from Code import gini, Calibration, PSI, Summary, Data
import Summary, gini, Calibration, PSI, Data, Customization


# Construct the path to the image
image_path = os.path.join(os.path.dirname(__file__), 'Images', 'NIMBUS_logo.png')
image = Image.open(image_path)

#Setting wide layout for streamlit app
st.set_page_config(
    page_title="Model Monitoring Dashboard", layout="wide", page_icon=image
)

# #CSS to hide the Streamlit stop button
# hide_st_button = """
#     <style>
#     #MainMenu {visibility: hidden;}
#     footer {visibility: hidden;}
#     header {visibility: hidden;}
#     .stApp > header {visibility: hidden;}
#     </style>
# """
# st.markdown(hide_st_button, unsafe_allow_html=True)

# Initial change log data
initial_change_log_data = [
    {"Date": "201904", "Description": "Validation completed"},
    {"Date": "201906", "Description": "1st Monitoring Report (Quarterly Updated)"},
    {"Date": "201909", "Description": "Implemented in System"},
    {"Date": "201909", "Description": "Model re-calibration"},
    {"Date": "202211", "Description": "Last Validation"},
    {"Date": "202212", "Description": "Go Live"},
    {"Date": "202305", "Description": "Tier 1"},
    {"Date": "202306", "Description": "Implemented in System (IFRS9 version)"},
    {"Date": "202308", "Description": "Tier 1"},
    {"Date": "202310", "Description": "Monitoring aligned with IFRS9 implementation"},
    {"Date": "202311", "Description": "Tier 1"},
    {"Date": "202404", "Description": "Tier 1"},
    {"Date": "202404", "Description": "Tier 2"},
]

if 'df_change_log' not in st.session_state:
    st.session_state.df_change_log = pd.DataFrame(initial_change_log_data)
    
# Function to load a PowerPoint presentation from BytesIO
def load_presentation_from_bytesio(presentation_bytesio):
    return Presentation(presentation_bytesio)

# Function to apply title styling
def style_title(title_shape):
    title_shape.left = Inches(0.5)
    title_shape.top = Inches(0.2)
    title_shape.width = Inches(9)
    title_shape.height = Inches(0.5)
    title_text_frame = title_shape.text_frame
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(22)
            # Check if color is in session state; if not, use default color
            if 'title_font_color' in st.session_state:
                run.font.color.rgb = RGBColor.from_string(st.session_state.title_font_color[1:])
            else:
                run.font.color.rgb = RGBColor(0, 0, 0)  # Default color black
            
            
# Function to apply title styling and background color
def set_slide_background_and_title_style(slide, title_shape, slide_index):
    
    # Create presentation object
    prs = Presentation()
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    # Check if color is in session state; if not, use default color
    if 'bg_color' in st.session_state:
        fill.fore_color.rgb = RGBColor.from_string(st.session_state.bg_color[1:])
    else:
        fill.fore_color.rgb = RGBColor(0x0A, 0x3D, 0x62)
    
    #Add slide number to title slide
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    ribbon_height = Inches(0.28)  # Adjust height as needed
    half_ribbon_width = slide_width / 2
    ribbon_top = slide_height - ribbon_height
    
    # Add the second half of the ribbon (dark blue)
    ribbon_left = half_ribbon_width  # Adjust the left position for the second half
    shape_blue = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, ribbon_left, ribbon_top, half_ribbon_width, ribbon_height
    )
    fill_blue = shape_blue.fill
    fill_blue.solid()
    # Check if olor is in session state; if not, use default color
    if 'bg_color' in st.session_state:
        fill_blue.fore_color.rgb = RGBColor.from_string(st.session_state.bg_color[1:])  # Dark blue color
    else:
        fill_blue.fore_color.rgb = RGBColor(0x0A, 0x3D, 0x62) # Dark blue color
    line_blue = shape_blue.line
    # Check if title_font_color is in session state; if not, use default color black
    if 'bg_color' in st.session_state:
        line_blue.color.rgb = RGBColor.from_string(st.session_state.bg_color[1:]) # Dark blue outline color
    else:
        line_blue.color.rgb = RGBColor(0x0A, 0x3D, 0x62) # Dark blue outline color
    
    text_frame = shape_blue.text_frame
    text_frame.text = f" {slide_index}\t"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(9)
    # Check if color is in session state; if not, use default color
    if 'ribbon_font_color' in st.session_state:
        p.font.color.rgb = RGBColor.from_string(st.session_state.ribbon_font_color[1:])  # White text
    else:
        p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    p.alignment = PP_ALIGN.RIGHT  # Align text to the right

    # Apply title styling
    title_shape.left = Inches(1)
    title_shape.top = Inches(2.6)  # Centered vertically
    title_shape.width = Inches(8)
    title_shape.height = Inches(2)
    title_text_frame = title_shape.text_frame
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(42)
            run.font.bold = True
            # Check if color is in session state; if not, use default color
            if 'font_color' in st.session_state:
                run.font.color.rgb = RGBColor.from_string(st.session_state.font_color[1:])  # White font color
            else:
                run.font.color.rgb = RGBColor(255, 255, 255)  # White font color
            
            
# Function to clone shapes (including GraphicFrame and pictures)
def clone_shape(original_shape, merged_slide):
    if original_shape.shape_type == 14:  # Placeholder
        ph = original_shape.placeholder_format
        merged_slide.placeholders[ph.idx].text = original_shape.text
    elif original_shape.shape_type == 13:  # Picture
        image_stream = BytesIO(original_shape.image.blob)
        left = original_shape.left
        top = original_shape.top
        width = original_shape.width
        height = original_shape.height
        merged_slide.shapes.add_picture(image_stream, left, top, width, height)
    else:  # Other shapes (including GraphicFrame)
        merged_shape = original_shape._element
        merged_slide.shapes._spTree.append(merged_shape)
    
# Function to add a slide with the ribbon and logo
def ppt_ribbon_and_logo(slide, slide_index):
    # Create presentation object
    prs = Presentation()
    # Determine slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Calculate ribbon dimensions
    ribbon_height = Inches(0.28)  # Adjust height as needed
    half_ribbon_width = slide_width / 2

    # Calculate positions for ribbons
    ribbon_left = Inches(0)
    ribbon_top = slide_height - ribbon_height

    # Add the first half of the ribbon (amber)
    shape_amber = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, ribbon_left, ribbon_top, half_ribbon_width, ribbon_height
    )
    fill_amber = shape_amber.fill
    fill_amber.solid()
    # Check if color is in session state; if not, use default color
    if 'ribbon_color_1' in st.session_state:
        fill_amber.fore_color.rgb = RGBColor.from_string(st.session_state.ribbon_color_1[1:])  # Amber color
    else:
        fill_amber.fore_color.rgb = RGBColor(255, 191, 0)  # Amber color

    line_amber = shape_amber.line
    # Check if color is in session state; if not, use default color
    if 'ribbon_color_1' in st.session_state:
        line_amber.color.rgb = RGBColor.from_string(st.session_state.ribbon_color_1[1:])  # Amber outline color
    else:
        line_amber.color.rgb = RGBColor(255, 191, 0)  # Amber color

    # Add text to the amber ribbon
    text_frame_amber = shape_amber.text_frame
    text_frame_amber.text = "\tModel Monitoring"
    text_frame_amber.paragraphs[0].font.size = Pt(9)
    # Check if title_font_color is in session state; if not, use default color black
    if 'ribbon_font_1' in st.session_state:
        text_frame_amber.paragraphs[0].font.color.rgb = RGBColor.from_string(st.session_state.ribbon_font_1[1:])  # Black text
    else:
        text_frame_amber.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text

    # Add the second half of the ribbon (dark blue)
    ribbon_left = half_ribbon_width  # Adjust the left position for the second half
    shape_blue = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, ribbon_left, ribbon_top, half_ribbon_width, ribbon_height
    )
    fill_blue = shape_blue.fill
    fill_blue.solid()
    # Check if olor is in session state; if not, use default color
    if 'ribbon_color_2' in st.session_state:
        fill_blue.fore_color.rgb = RGBColor.from_string(st.session_state.ribbon_color_2[1:])  # Dark blue color
    else:
        fill_blue.fore_color.rgb = RGBColor(255, 191, 0)   # Dark blue color

    line_blue = shape_blue.line
    # Check if title_font_color is in session state; if not, use default color black
    if 'ribbon_color_2' in st.session_state:
        line_blue.color.rgb = RGBColor.from_string(st.session_state.ribbon_color_2[1:]) # Dark blue outline color
    else:
        line_blue.color.rgb = RGBColor(255, 191, 0)  # Dark blue outline color
    
    # Add slide number
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    ribbon_height = Inches(0.28)  # Adjust height as needed
    half_ribbon_width = slide_width / 2
    ribbon_top = slide_height - ribbon_height
    
    text_frame = shape_blue.text_frame
    text_frame.text = f" {slide_index}\t"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(9)
    # Check if color is in session state; if not, use default color
    if 'ribbon_font_2' in st.session_state:
        p.font.color.rgb = RGBColor.from_string(st.session_state.ribbon_font_2[1:])  # White text
    else:
        p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    p.alignment = PP_ALIGN.RIGHT  # Align text to the right
    
    # Add logo image to the upper right corner
    base_dir = os.path.dirname(__file__)
    
    # Logo path
    logo_path = os.path.join(base_dir, 'Images', 'ENBD_s.jpg')  # Replace with your actual file path
    logo_left = slide_width - Inches(0.6)  # Adjust position as needed
    logo_top = Inches(0.15)  # Adjust position as needed
    logo_height = Inches(0.45)  # Adjust size as needed
    slide.shapes.add_picture(logo_path, logo_left, logo_top, height=logo_height)

# Function to create a PowerPoint presentation for Change Log
def create_change_log_presentation(df_change_log):
    prs = Presentation()
    
    slide_index = 1  # To keep track of the slide index

    # Add a slide with title and content layout
    slide_layout = prs.slide_layouts[5]  # Use the blank layout
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    ppt_ribbon_and_logo(slide, slide_index)

    # Title slide
    title_shape = shapes.title
    title_shape.text = "Change Log"
    title_shape.left = Inches(0.5)
    title_shape.top = Inches(0.2)
    title_shape.width = Inches(9)
    title_shape.height = Inches(0.5)
    title_text_frame = title_shape.text_frame
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(22)
            # Check if color is in session state; if not, use default color
            if 'title_font_color' in st.session_state:
                run.font.color.rgb = RGBColor.from_string(st.session_state.title_font_color[1:])
            else:
                run.font.color.rgb = RGBColor(0, 0, 0)

    # Determine table size and position
    slide_width = prs.slide_width
    table_width = Inches(6)  # Adjust width to fit content
    table_height = Inches(5)  # Adjusted to fit the table within the slide
    left = (slide_width - table_width) / 2  # Center the table horizontally
    top = Inches(1.1)  # Center the table vertically

    rows, cols = df_change_log.shape
    table_shape = shapes.add_table(rows + 1, cols, left, top, table_width, table_height)
    table = table_shape.table
    table_style_id = table._tbl.tblPr.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}tableStyleId"
    )
    table_style_id.text = "{5940675A-B579-460E-94D1-54222C63F5DA}"

    # Set column names and font size
    for col_idx, col_name in enumerate(df_change_log.columns):
        cell = table.cell(0, col_idx)
        cell.text = col_name
        cell.fill.solid()
        # Check if color is in session state; if not, use default color
        if 'row_bg_color' in st.session_state:
            cell.fill.fore_color.rgb = RGBColor.from_string(st.session_state.row_bg_color[1:]) # Teal
        else:
            cell.fill.fore_color.rgb = RGBColor(0x00, 0x80, 0x80) # Teal
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(11.5)
                # Check if color is in session state; if not, use default color
                if 'row_font_color' in st.session_state:
                    run.font.color.rgb = RGBColor.from_string(st.session_state.row_font_color[1:])  # White font color
                else:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                

    # Adding data to table with font size set to 10
    for i in range(len(df_change_log)):
        cell1 = table.cell(i + 1, 0)
        cell1.text = str(df_change_log.iloc[i, 0])
        for paragraph in cell1.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(11)
                # Check if color is in session state; if not, use default color
                if 'content_font_color' in st.session_state:
                    run.font.color.rgb = RGBColor.from_string(st.session_state.content_font_color[1:])
                else:
                    run.font.color.rgb = RGBColor(0, 0, 0)
                

        cell2 = table.cell(i + 1, 1)
        cell2.text = str(df_change_log.iloc[i, 1])
        for paragraph in cell2.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(11)
                # Check if color is in session state; if not, use default color
                if 'content_font_color' in st.session_state:
                    run.font.color.rgb = RGBColor.from_string(st.session_state.content_font_color[1:])
                else:
                    run.font.color.rgb = RGBColor(0, 0, 0)

    # Set row heights
    for row in range(rows + 1):
        table.rows[row].height = Inches(0.3)

    return prs

# Function to create a download link for the Change Log presentation
def create_download_link_for_change_log_ppt(df_change_log):
    presentation = create_change_log_presentation(df_change_log)
    ppt_stream = BytesIO()
    presentation.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream
    
#Function to create powerpoint presentation for overview  
def create_ppt_overview_image(image_path):
    try:
        # Create a presentation object
        prs = Presentation()
        
        slide_index = 1  # To keep track of the slide index
        
        # Add a slide with title and content layout
        slide_layout = prs.slide_layouts[5]  # Use the blank layout
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        
        ppt_ribbon_and_logo(slide, slide_index)
        
        # Add and style title with a shorter height
        title_shape = shapes.title
        title_shape.text = "Overview"
        title_shape.left = Inches(0.5)
        title_shape.top = Inches(0.3)
        title_shape.width = Inches(9)
        title_shape.height = Inches(0.5)
        title_text_frame = title_shape.text_frame
        for paragraph in title_text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(22)
                # Check if color is in session state; if not, use default color
                if 'title_font_color' in st.session_state:
                    run.font.color.rgb = RGBColor.from_string(st.session_state.title_font_color[1:])
                else:
                    run.font.color.rgb = RGBColor(0, 0, 0)
        
        # Open the image
        image = Image.open(image_path)
        
        # Save the image to a byte stream
        img_byte_arr = io.BytesIO()
        image.save(img_byte_arr, format='JPEG')
        img_byte_arr.seek(0)
        
        # Add the image to the slide
        slide.shapes.add_picture(img_byte_arr, Inches(0.2), Inches(1), Inches(9.6), Inches(5.3))
        
        # Save the presentation to a byte stream
        ppt_byte_arr = io.BytesIO()
        prs.save(ppt_byte_arr)
        ppt_byte_arr.seek(0)
        
        return ppt_byte_arr
    except Exception as e:
        st.error(f"Error creating PowerPoint: {e}")
        return None
    
#Function to create download link for overview module
def create_ppt_download_overview(image_path):
    ppt_byte_arr = create_ppt_overview_image(image_path)
    return ppt_byte_arr

# Inject custom styles at the beginning of the Streamlit app
def inject_custom_styles():
    style = """
    <style>
    .stButton button,
    .stDownloadButton button {
        border: 1px solid #ccc !important;
        color: #333 !important;
        background-color: #f0f0f0 !important;
        box-shadow: 3px 3px 5px #aaa !important;
        padding: 0.5em 1em;
        font-size: 1em;
        border-radius: 5px;
        transition: transform 0.1s ease-in-out;
    }
    .stButton button:hover,
    .stDownloadButton button:hover {
        transform: translateY(-2px);
        box-shadow: 5px 5px 7px #999 !important;
    }
    .stButton button:active,
    .stDownloadButton button:active {
        transform: translateY(2px);
        box-shadow: 1px 1px 2px #bbb !important;
    }
    .ppt-download-button,
    .excel-download-button {
        display: inline-flex;
        align-items: center;
        border: 1px solid #ccc !important;
        color: #333 !important;
        background-color: #f0f0f0 !important;
        box-shadow: 3px 3px 5px #aaa !important;
        padding: 0.5em 1em;
        font-size: 1em;
        border-radius: 5px;
        transition: transform 0.1s ease-in-out;
    }
    .ppt-download-button:hover,
    .excel-download-button:hover {
        transform: translateY(-2px);
        box-shadow: 5px 5px 7px #999 !important;
    }
    .ppt-download-button:active,
    .excel-download-button:active {
        transform: translateY(2px);
        box-shadow: 1px 1px 2px #bbb !important;
    }
    .ppt-download-button img,
    .excel-download-button img {
        margin-left: -0.1em;
        margin-right: -0.1em;
        width: 32px;
        height: 27px;
    }
    </style>
    """
    st.markdown(style, unsafe_allow_html=True)

# Inject styles
inject_custom_styles()

# Function to merge four PowerPoint presentations and apply title styling to specific slides
def merge_presentations(presentation1, presentation2):
    merged_presentation = Presentation()
    slide_index = 1  # To keep track of the slide index

    # Copy slides from presentation1 (Overview) first
    for slide in presentation1.slides:
        slide_layout = merged_presentation.slide_master.slide_layouts[5]
        merged_slide = merged_presentation.slides.add_slide(slide_layout)

        for shape in slide.shapes:
            clone_shape(shape, merged_slide)

        # Apply title styling if there is a title and the slide index is 1
        if merged_slide.shapes.title:
            if slide_index in [1]:
                style_title(merged_slide.shapes.title)
                ppt_ribbon_and_logo(merged_slide, slide_index)
        
        slide_index += 1
        
    # Copy slides from presentation2 (Change Log)
    for slide in presentation2.slides:
        slide_layout = merged_presentation.slide_master.slide_layouts[5]
        merged_slide = merged_presentation.slides.add_slide(slide_layout)

        for shape in slide.shapes:
            clone_shape(shape, merged_slide)

        # Apply title styling if there is a title and the slide index is 1
        if merged_slide.shapes.title:
            if slide_index in [2]:
                style_title(merged_slide.shapes.title)
                ppt_ribbon_and_logo(merged_slide, slide_index)
        
        slide_index += 1
        
    return merged_presentation


class MultiApp:
    def __init__(self):
        self.apps = []

    def add_app(self, title, func):
        self.apps.append({
            "title": title,
            "function": func
        })

    def run(self):
        with st.sidebar:
            app = option_menu(
                menu_title='PL - Scorecard Model Monitoring',
                options=['Overview', 'Change Log', 'Summary', 'Risk Metrics', 'Data', 'PPT Customization'],
                icons=['cast', 'book', 'file-text', 'exclamation-triangle', 'database', 'palette'],
                menu_icon='bank',
                default_index=0,
                styles={
                    "container": {"padding": "20px", "background-color": 'white', "border-radius": "5px"},
                    "icon": {"color": "#008080", "font-size": "17px"},
                    "nav-link": {"color": "black", "font-size": "14px", "text-align": "left", "margin": "2px", "--hover-color": "#B9DFFE", "font-family": "sans-serif"},
                    "nav-link-selected": {"background-color": "#AFDBF5", "border-radius": "5px"},
                }
            )
            
        # Main content area
        if app == 'Risk Metrics':
            # Display metric selection in the main content area
            risk_metric = st.sidebar.selectbox(
                'Select Metric',
                options=['Gini', 'Calibration', 'PSI'],
                index=0
            )
            app = risk_metric

            # # Display the appropriate module based on the selected metric
            # if risk_metric == 'Gini':
            #     gini.app()
            # elif risk_metric == 'Calibration':
            #     Calibration.app()
            # elif risk_metric == 'PSI':
            #     PSI.app()


        if app == "Overview":
            # Centered text in the floating container
            st.markdown(
                """
                <div style='background-color: #008080; border-radius: 5px; padding: 1px;'>
                    <h1 style='text-align: center; font-size: 28px; color: white;'>Welcome to PL - Scorecard Model Monitoring Dashboard</h1>
                </div>
                """,
                unsafe_allow_html=True
            )
            
            # Define custom CSS to center text and set font size
            custom_css = """
            <style>
            .centered-overview {
                text-align: center;
                font-size: 20px;
            }
            </style>
            """
            st.markdown(custom_css, unsafe_allow_html=True)   
            
            # Use the custom CSS class to center and resize the text
            st.markdown("<div class='centered-overview'><b>Overview</b></div>", unsafe_allow_html=True)
            
            # Display image below the text
            image_path = os.path.join(os.path.dirname(__file__), 'Images', 'Overview_1.jpg')
            image = Image.open(image_path)
            st.image(image, use_column_width=True)
            
            # Base directory of the current script
            base_dir = os.path.dirname(__file__)
            
            # Construct the path for the image and use it in the create_ppt_download_overview function
            ppt_data_overview = create_ppt_download_overview(
                image_path=os.path.join(base_dir, 'Images', 'Overview_1.jpg')
            )
            
            st.session_state.ppt_data_overview = ppt_data_overview
                
            # Base directory of the current script
            base_dir = os.path.dirname(__file__)
            
            # Construct the path for the image
            image_path = os.path.join(base_dir, "Images", "ppt_logo.png")
            
            # Read and encode the image to base64
            with open(image_path, "rb") as image_file:
                image_base64 = base64.b64encode(image_file.read()).decode()
            
            # Convert BytesIO to bytes
            ppt_data_overview_bytes = st.session_state.ppt_data_overview.getvalue()
            
            # Encode the PPT data to base64
            ppt_data_overview_base64 = base64.b64encode(ppt_data_overview_bytes).decode()
            
            # Help text
            help_text = "Click here to download the PowerPoint presentation"
            
            # Create the HTML for the button with image and help text
            button_html = f"""
            <a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{ppt_data_overview_base64}" download="Overview.pptx" class="ppt-download-button" title="{help_text}">
                <img src="data:image/png;base64,{image_base64}" alt="Download PPT">
            </a>
            """
            
            # Display the custom download button in Streamlit
            st.markdown(button_html, unsafe_allow_html=True)

# =============================================================================
#             #Creating new consolidated excel download button where we can download the excel template as there's
# =============================================================================
   
            
            #Threshods Gini
            def threshold_selection_gini():
                # Base directory of the current script
                base_dir = os.path.dirname(__file__)
            
                # Construct the path for the pickle file
                pickle_path = os.path.join(base_dir, 'pkl', 'model_gini.pkl')
            
                try:
                    with open(pickle_path, 'rb') as f:
                        thresholds_gini = pickle.load(f)
                except FileNotFoundError:
                    thresholds_gini = None
            
                if thresholds_gini:
                    green_threshold = thresholds_gini['green_gini']['value']
                    amber_lower = thresholds_gini['amber_gini']['lower']
                    amber_upper = thresholds_gini['amber_gini']['upper']
                    red_threshold = thresholds_gini['red_gini']['value']
                else:
                    # Initialize with default values if no thresholds are found
                    green_threshold = None
                    amber_lower = None
                    amber_upper = None
                    red_threshold = None
                
                return green_threshold, amber_lower, amber_upper, red_threshold
            
            # Get thresholds
            green_threshold, amber_lower, amber_upper, red_threshold = threshold_selection_gini()
            
            # Function to create the threshold selection UI
            def threshold_selection_calibration():
                # Base directory of the current script
                base_dir = os.path.dirname(__file__)
            
                # Construct the path for the pickle file
                pickle_path = os.path.join(base_dir, 'pkl', 'model_calibration.pkl')
            
                try:
                    with open(pickle_path, 'rb') as j:
                        thresholds_calibration = pickle.load(j)
                except FileNotFoundError:
                    thresholds_calibration = None

                if thresholds_calibration:
                    green_threshold_1 = thresholds_calibration.get('green_calibration_1', {}).get('value', -0.075)
                    amber_lower_1 = thresholds_calibration.get('amber_calibration_1', {}).get('lower', -0.15)
                    amber_upper_1 = thresholds_calibration.get('amber_calibration_1', {}).get('upper', -0.075)
                    red_threshold_1 = thresholds_calibration.get('red_calibration_1', {}).get('value', -0.15)
                    green_threshold_2 = thresholds_calibration.get('green_calibration_2', {}).get('value', 0.20)
                    amber_lower_2 = thresholds_calibration.get('amber_calibration_2', {}).get('lower', 0.20)
                    amber_upper_2 = thresholds_calibration.get('amber_calibration_2', {}).get('upper', 0.30)
                    red_threshold_2 = thresholds_calibration.get('red_calibration_2', {}).get('value', 0.30)
                else:
                    # Default values when thresholds are not provided or loaded
                    green_threshold_1 = None
                    amber_lower_1 = None
                    amber_upper_1 = None
                    red_threshold_1 = None
                    green_threshold_2 = None
                    amber_lower_2 = None
                    amber_upper_2 = None
                    red_threshold_2 = None
                
                return green_threshold_1, amber_lower_1, amber_upper_1, red_threshold_1, green_threshold_2, amber_lower_2, amber_upper_2, red_threshold_2
            
            # Get thresholds
            green_threshold_1, amber_lower_1, amber_upper_1, red_threshold_1, green_threshold_2, amber_lower_2, amber_upper_2, red_threshold_2 = threshold_selection_calibration()
            
            # Function to create the threshold selection UI for PSI
            def threshold_selection_PSI():
                # Base directory of the current script
                base_dir = os.path.dirname(__file__)
            
                # Construct the path for the pickle file
                pickle_path = os.path.join(base_dir, 'pkl', 'model_psi.pkl')
            
                try:
                    with open(pickle_path, 'rb') as k:
                        thresholds_psi = pickle.load(k)
                except FileNotFoundError:
                    thresholds_psi = None
                    
                if thresholds_psi:
                    green_threshold_psi = thresholds_psi['green_psi']['value']
                    amber_lower_psi = thresholds_psi['amber_psi']['lower']
                    amber_upper_psi = thresholds_psi['amber_psi']['upper']
                    red_threshold_psi = thresholds_psi['red_psi']['value']
                else:
                    # Default values when UI is not shown
                    green_threshold_psi = None
                    amber_lower_psi =None
                    amber_upper_psi = None
                    red_threshold_psi = None
                    
                return green_threshold_psi, amber_lower_psi, amber_upper_psi, red_threshold_psi
            
            # Get thresholds
            green_threshold_psi, amber_lower_psi, amber_upper_psi, red_threshold_psi = threshold_selection_PSI()
            
            #Uploading template with the same customization as ENBD's
            # Base directory of the current script
            base_dir = os.path.dirname(__file__)
            
            # File paths
            existing_file_path = os.path.join(base_dir, 'Datasets', 'Excel_workbook.xlsx')
            new_data_file_path = os.path.join(base_dir, 'Datasets', 'Support_2.xlsx')
            sheet_name = 'support'
            
            
            # Define fill colors for conditional formatting
            fill_red = PatternFill(start_color='FF0000', end_color='FF0000', fill_type='solid')
            fill_amber = PatternFill(start_color='FFC107', end_color='FFC107', fill_type='solid')
            fill_green = PatternFill(start_color='00B050', end_color='00B050', fill_type='solid')

            # Check if the existing file exists
            def update_support_sheet(existing_file_path, new_data_file_path, sheet_name):
                # Check if the existing file exists
                if not os.path.exists(existing_file_path):
                    return f"Error: The existing template {existing_file_path} does not exist."
                
                # Check if the new data file exists
                if not os.path.exists(new_data_file_path):
                    return "Error: The Data recieved from Nimbus 'new_data_file_path' does not exist."
            
                try:
                    # Load the new data file and check if the sheet exists
                    new_data_book = load_workbook(new_data_file_path, read_only=True)
                    if sheet_name not in new_data_book.sheetnames:
                        return f"Error: The sheet named '{sheet_name}' does not exist in 'new_data_file_path'."
                    
                    # Load the new data into a pandas DataFrame
                    new_data_df = pd.read_excel(new_data_file_path, sheet_name=sheet_name)
            
                    # Load existing workbook using openpyxl
                    book = load_workbook(existing_file_path)
            
                    # Access the existing "Support" sheet
                    if sheet_name in book.sheetnames:
                        support_sheet = book[sheet_name]
                    else:
                        return f"Error: The sheet named '{sheet_name}' does not exist in {existing_file_path}."
            
                    # Clear the existing "Support" sheet
                    support_sheet.delete_rows(1, support_sheet.max_row)
            
                    # Write the new data into the "Support" sheet
                    for r in dataframe_to_rows(new_data_df, index=False, header=True):
                        support_sheet.append(r)
                        
                    # Apply highlighting to cells K25 and K40 if the sheet is "GINI"
                    if "3. Gini" in book.sheetnames:
                        gini_sheet = book["3. Gini"]
                    for cell in ['K25', 'K40']:
                        gini_sheet.conditional_formatting.add(
                            cell,
                            FormulaRule(
                                formula=[f'{cell}<={red_threshold}'],
                                stopIfTrue=True,
                                fill=fill_red
                            )
                        )
                        gini_sheet.conditional_formatting.add(
                            cell,
                            FormulaRule(
                                formula=[f'AND({cell}>{amber_lower}, {cell}<={amber_upper})'],
                                stopIfTrue=True,
                                fill=fill_amber
                            )
                        )
                        gini_sheet.conditional_formatting.add(
                            cell,
                            FormulaRule(
                                formula=[f'{cell}>{green_threshold}'],
                                stopIfTrue=True,
                                fill=fill_green
                            )
                        )
                        
                    # Apply highlighting to cells L23 and L38 if the sheet is "Calibration"
                    if "4. Calibration" in book.sheetnames:
                        Cal_sheet = book["4. Calibration"]
                    for cell in ['L23', 'L38']:
                        # Highlight in red if value is less than or equal to red_threshold and less than 0
                        Cal_sheet.conditional_formatting.add(
                            cell,
                            FormulaRule(
                                formula=[f'AND({cell}<0, {cell}<{red_threshold_1})'],
                                stopIfTrue=True,
                                fill=fill_red
                            )
                        )
                        # Highlight in amber if value is between amber_lower and amber_upper and less than 0
                        Cal_sheet.conditional_formatting.add(
                            cell,
                            FormulaRule(
                                formula=[f'AND({cell}<0, {cell}>={amber_lower_1}, {cell}<={amber_upper_1})'],
                                stopIfTrue=True,
                                fill=fill_amber
                            )
                        )
                        # Highlight in green if value is greater than green_threshold and less than 0
                        Cal_sheet.conditional_formatting.add(
                            cell,
                            FormulaRule(
                                formula=[f'AND({cell}<0, {cell}>{green_threshold_1})'],
                                stopIfTrue=True,
                                fill=fill_green
                            )
                        )
                        # Highlight in red if value is less than or equal to red_threshold and greater than or equal to 0
                        Cal_sheet.conditional_formatting.add(
                            cell,
                            FormulaRule(
                                formula=[f'AND({cell}>=0, {cell}>{red_threshold_2})'],
                                stopIfTrue=True,
                                fill=fill_red
                            )
                        )
                        # Highlight in amber if value is between amber_lower and amber_upper and greater than or equal to 0
                        Cal_sheet.conditional_formatting.add(
                            cell,
                            FormulaRule(
                                formula=[f'AND({cell}>=0, {cell}>={amber_lower_2}, {cell}<={amber_upper_2})'],
                                stopIfTrue=True,
                                fill=fill_amber
                            )
                        )
                        # Highlight in green if value is greater than green_threshold and greater than or equal to 0
                        Cal_sheet.conditional_formatting.add(
                            cell,
                            FormulaRule(
                                formula=[f'AND({cell}>=0, {cell}<{green_threshold_2})'],
                                stopIfTrue=True,
                                fill=fill_green
                            )
                        )
                        
                    # Apply highlighting to cells of row 36 if the sheet is "PSI"    
                    if "5. PSI" in book.sheetnames:
                        psi_sheet = book["5. PSI"]
                    for cell in ['D36', 'E36', 'F36', 'G36', 'H36', 'I36', 'J36', 'K36', 'L36', 'M36', 'N36', 'O36', 'P36']:
                        # Highlight in red if value is less than or equal to red_threshold and less than 0
                        psi_sheet.conditional_formatting.add(
                            cell,
                            FormulaRule(
                                formula=[f'{cell}>{red_threshold_psi}'],
                                stopIfTrue=True,
                                fill=fill_red
                            )
                        )
                        # Highlight in amber if value is between amber_lower and amber_upper and less than 0
                        psi_sheet.conditional_formatting.add(
                            cell,
                            FormulaRule(
                                formula=[f'AND({cell}>{amber_lower_psi}, {cell}<={amber_upper_psi})'],
                                stopIfTrue=True,
                                fill=fill_amber
                            )
                        )
                        # Highlight in green if value is greater than green_threshold and less than 0
                        psi_sheet.conditional_formatting.add(
                            cell,
                            FormulaRule(
                                formula=[f'{cell}<={green_threshold_psi}'],
                                stopIfTrue=True,
                                fill=fill_green
                            )
                        )
                        
                        # Load the workbook and select the summary sheet
                        if "2. Summary" in book.sheetnames:
                            summary_sheet = book["2. Summary"]
                        
                        # # Apply formula-based conditional formatting to cells K25 and K40
                        # for cell in ['D15']:
                        #     summary_sheet.conditional_formatting.add(
                        #         cell,
                        #         FormulaRule(
                        #             formula=[f'{cell}<={red_threshold}'],
                        #             stopIfTrue=True,
                        #             fill=fill_red
                        #         )
                        #     )
                        #     summary_sheet.conditional_formatting.add(
                        #         cell,
                        #         FormulaRule(
                        #             formula=[f'AND({cell}>{amber_lower}, {cell}<={amber_upper})'],
                        #             stopIfTrue=True,
                        #             fill= fill_amber
                        #         )
                        #     )
                        #     summary_sheet.conditional_formatting.add(
                        #         cell,
                        #         FormulaRule(
                        #             formula=[f'{cell}>{green_threshold}'],
                        #             stopIfTrue=True,
                        #             fill=fill_green
                        #         )
                        #     )
                            
                            # Set formula in cell E15 based on the value of cell D15
                            summary_sheet["E15"].value = (
                                f'=IF(D15>{green_threshold}, 1, IF(AND(D15>{amber_lower}, D15<={amber_upper}), 0, IF(D15<={red_threshold}, -1, "")))'
                            )

                            # Apply icon set conditional formatting using the thresholds variable
                            icon_set = IconSet(
                                iconSet='3Symbols2',
                                cfvo=[FormatObject(type='num', val=-1),
                                      FormatObject(type='num', val=0),
                                      FormatObject(type='num', val=1)],
                                showValue=None,  # This can be True or False to show/hide values
                                reverse=False,
                            )
                            dxf = DifferentialStyle()
                            rule = Rule(type='iconSet', dxf=dxf, iconSet=icon_set)
                            summary_sheet.conditional_formatting.add('E15', rule)
                                                    
                        # # Apply formula-based conditional formatting to cells L23 and L38
                        # for cell in ['F15']:
                        #     summary_sheet.conditional_formatting.add(
                        #         cell,
                        #         FormulaRule(
                        #             formula=[f'AND({cell}<0, {cell}<{red_threshold_1})'],
                        #             stopIfTrue=True,
                        #             fill=fill_red
                        #         )
                        #     )
                        #     summary_sheet.conditional_formatting.add(
                        #         cell,
                        #         FormulaRule(
                        #             formula=[f'AND({cell}<0, {cell}>={amber_lower_1}, {cell}<={amber_upper_1})'],
                        #             stopIfTrue=True,
                        #             fill=fill_amber
                        #         )
                        #     )
                        #     summary_sheet.conditional_formatting.add(
                        #         cell,
                        #         FormulaRule(
                        #             formula=[f'AND({cell}<0, {cell}>{green_threshold_1})'],
                        #             stopIfTrue=True,
                        #             fill=fill_green
                        #         )
                        #     )
                        #     summary_sheet.conditional_formatting.add(
                        #         cell,
                        #         FormulaRule(
                        #             formula=[f'AND({cell}>=0, {cell}>{red_threshold_2})'],
                        #             stopIfTrue=True,
                        #             fill=fill_red
                        #         )
                        #     )
                        #     summary_sheet.conditional_formatting.add(
                        #         cell,
                        #         FormulaRule(
                        #             formula=[f'AND({cell}>=0, {cell}>={amber_lower_2}, {cell}<={amber_upper_2})'],
                        #             stopIfTrue=True,
                        #             fill=fill_amber
                        #         )
                        #     )
                        #     summary_sheet.conditional_formatting.add(
                        #         cell,
                        #         FormulaRule(
                        #             formula=[f'AND({cell}>=0, {cell}<{green_threshold_2})'],
                        #             stopIfTrue=True,
                        #             fill=fill_green
                        #         )
                        #     )
                            
                            # Set formula in cell G15 based on the value of cell F15
                            summary_sheet["G15"].value = (
                                f'=IF(AND(F15<{red_threshold_1},F15<0), -1, '
                                f'IF(AND(F15>={amber_lower_1}, F15<={amber_upper_1}), 0, '
                                f'IF(AND(F15>{green_threshold_1},F15<0), 1, '
                                f'IF(AND(F15>{red_threshold_2}, F15>=0), -1, '
                                f'IF(AND(F15>={amber_lower_2}, F15<={amber_upper_2}), 0, '
                                f'IF(AND(F15<{green_threshold_2},F15>=0), 1, ""))))))'
                            )
                            
                            # Apply icon set conditional formatting using the thresholds variable
                            icon_set = IconSet(
                                iconSet='3Symbols2',
                                cfvo=[FormatObject(type='num', val=-1),
                                      FormatObject(type='num', val=0),
                                      FormatObject(type='num', val=1)],
                                showValue=None,  # This can be True or False to show/hide values
                                reverse=False,
                            )
                            dxf = DifferentialStyle()
                            rule = Rule(type='iconSet', dxf=dxf, iconSet=icon_set)
                            summary_sheet.conditional_formatting.add('G15', rule)
                            
                        
                        # # Apply formula-based conditional formatting to cells of row 36
                        # for cell in ['H15']:
                        #     summary_sheet.conditional_formatting.add(
                        #         cell,
                        #         FormulaRule(
                        #             formula=[f'{cell}>{red_threshold_psi}'],
                        #             stopIfTrue=True,
                        #             fill=fill_red
                        #         )
                        #     )
                        #     summary_sheet.conditional_formatting.add(
                        #         cell,
                        #         FormulaRule(
                        #             formula=[f'AND({cell}>{amber_lower_psi}, {cell}<={amber_upper_psi})'],
                        #             stopIfTrue=True,
                        #             fill=fill_amber
                        #         )
                        #     )
                        #     summary_sheet.conditional_formatting.add(
                        #         cell,
                        #         FormulaRule(
                        #             formula=[f'{cell}<={green_threshold_psi}'],
                        #             stopIfTrue=True,
                        #             fill=fill_green
                        #         )
                        #     )
                            
                            # Set formula in cell G15 based on the value of cell F15
                            summary_sheet["I15"].value = (
                                f'=IF(H15>{red_threshold_psi}, -1, '
                                f'IF(AND(H15>{amber_lower_psi}, H15<={amber_upper_psi}), 0, '
                                f'IF(H15<={green_threshold_psi}, 1, "")))'
                            )
                            
                            # Apply icon set conditional formatting using the thresholds variable
                            icon_set = IconSet(
                                iconSet='3Symbols2',
                                cfvo=[FormatObject(type='num', val=-1),
                                      FormatObject(type='num', val=0),
                                      FormatObject(type='num', val=1)],
                                showValue=None,  # This can be True or False to show/hide values
                                reverse=False,
                            )
                            dxf = DifferentialStyle()
                            rule = Rule(type='iconSet', dxf=dxf, iconSet=icon_set)
                            summary_sheet.conditional_formatting.add('I15', rule)
                        
                    # Save the updated workbook to a BytesIO object
                    buffer = BytesIO()
                    book.save(buffer)
                    buffer.seek(0)
            
                    return buffer
                
            
                except Exception as e:
                    return f"Error: {e}"
            
            # Load the Excel data from the BytesIO object and extract the heading
            def load_excel_and_get_heading(excel_buffer):
                book = load_workbook(excel_buffer, data_only=True)
                support_sheet = book["support"]
                heading = support_sheet['B1'].value
                subheading = support_sheet['B2'].value
                monitoring = 'Monitoring  Report Date'
                date = support_sheet['O17'].value
                
                return heading, subheading, monitoring, date
                    
                    
            # Update the support sheet and get the BytesIO buffer with the updated file
            updated_file_buffer = update_support_sheet(existing_file_path, new_data_file_path, sheet_name)
            
            # Display the result or provide a download button if the update was successful
            if isinstance(updated_file_buffer, str) and "Error" in updated_file_buffer:
                st.error(updated_file_buffer)
            else:
                
                st.session_state.updated_workbook = updated_file_buffer
                heading, subheading, monitoring_date, date = load_excel_and_get_heading(new_data_file_path)
                
                # Save values to session state
                st.session_state.heading = heading
                st.session_state.subheading = subheading
                st.session_state.monitoring_date = monitoring_date
                st.session_state.date = date
                
                # Convert BytesIO to bytes
                excel_data_bytes = updated_file_buffer.getvalue()
            
                # Encode the Excel data to base64
                excel_data_base64 = base64.b64encode(excel_data_bytes).decode()
            
                # Base directory of the current script
                base_dir = os.path.dirname(__file__)
                
                # Construct the path for the image
                image_path = os.path.join(base_dir, "Images", "excel_logo.png")
                
                # Read and encode the image to base64
                with open(image_path, "rb") as image_file:
                    excel_image_base64 = base64.b64encode(image_file.read()).decode()
            
                # Help text for the Excel button
                help_text = "Click here to download the Consolidated Excel Workbook"
            
                # Create the HTML for the Excel download button with image and help text
                excel_button_html = f"""
                <a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{excel_data_base64}" download="Excel_template.xlsx" class="excel-download-button" title="{help_text}">
                    <img src="data:image/png;base64,{excel_image_base64}" alt="Download Excel">
                </a>
                """
                
                custom_css = """
                    <style>
                        .excel-download-button {
                            position: absolute;
                            top: -60px;
                            left: 95px;
                            cursor: pointer;
                            z-index: 10;
                        }
                    </style>
                    """
                st.markdown(custom_css, unsafe_allow_html=True)
                
                # Display the custom Excel download button in Streamlit
                st.markdown(excel_button_html, unsafe_allow_html=True)
                     

        elif app == "Change Log":
            st.markdown(
                """
                <div style='background-color: #008080; border-radius: 5px; padding: 1px;'>
                    <h1 style='text-align: center; font-size: 28px; color: white;'>PL- Scorecard Model Monitoring Change Log</h1>
                </div>
                """,
                unsafe_allow_html=True
            )
            
            # Define your CSS styles with custom colors
            css = """
                <style>
                    .custom-heading {
                        font-weight: bold;
                        font-size: 20px; /* Adjust size as needed */
                        color: #44546A; /* Custom color for the heading */
                        margin-bottom: 5px; /* Adjust spacing between headings */
                    }
                    .custom-subheading {
                        font-weight: bold;
                        font-size: 17px; /* Adjust size as needed */
                        color: #4472C4; /* Custom color for the subheading */
                        margin-bottom: 3px; /* Adjust spacing between subheading and date */
                    }
                    .custom-date {
                        font-weight: bold;
                        font-size: 16px; /* Adjust size as needed */
                    }
                </style>
            """
            st.markdown(css, unsafe_allow_html=True)
            
            # Initialize variables with default values
            heading = ""
            subheading = ""
            monitoring_date = ""
            date = ""
            
            
            #Calling variables from session_state for    
            if 'heading' in st.session_state:
                heading = st.session_state.heading
                
            if 'subheading' in st.session_state:
                subheading = st.session_state.subheading
            
            if 'monitoring_date' in st.session_state:
                monitoring_date = st.session_state.monitoring_date
            
            if 'date' in st.session_state:
                date = st.session_state.date
            
            # Display the extracted heading with custom CSS classes
            st.markdown(f"""
                <div class="custom-heading">{heading}</div>
                <div class="custom-date">{monitoring_date} {date}</div>
            """, unsafe_allow_html=True)
            
            #Showing dataframe in streamlit
            st.dataframe(st.session_state.df_change_log, width=1200)
            
            # Initialize session state for the flag
            if 'entry_added' not in st.session_state:
                st.session_state.entry_added = False
            
            #Adding Help section for button when we go to the button then it will show this help
            if st.button("Add New Entry", help="Click here to add an entry to the Change_log table"):
                st.session_state.show_form = True
        
            if st.session_state.get('show_form', False):
                with st.form("add_entry_form"):
                    new_date = st.text_input("Date")
                    new_description = st.text_input("Description")
                    submitted = st.form_submit_button("Add Entry")
        
                    if submitted:
                        if new_date and new_description:
                            new_entry = {"Date": new_date, "Description": new_description}
                            new_entry_df = pd.DataFrame([new_entry])
                            st.session_state.df_change_log = pd.concat([st.session_state.df_change_log, new_entry_df], ignore_index=True)
                            st.success("New entry added successfully!")
                            st.session_state.show_form = False
                            st.session_state.new_entry_df = new_entry_df  # Save new_entry_df in session state
                            st.session_state.entry_added = True  # Set the flag to indicate entry has been added
                            st.rerun()
                        else:
                            st.error("Please enter both Date and Description.")
                            
            # # Prepare Excel file in memory
            # excel_data = io.BytesIO()
            # with pd.ExcelWriter(excel_data, engine='openpyxl') as writer:
            #     st.session_state.df_change_log.to_excel(writer, index=False, sheet_name='Summary', startrow=12)
            # excel_data.seek(0)
            
            # # Convert BytesIO to bytes
            # excel_data_bytes = excel_data.getvalue()
            
            # # Encode the Excel data to base64
            # excel_data_base64 = base64.b64encode(excel_data_bytes).decode()
            
            # # Read and encode the Excel image to base64
            # with open("excel_logo.png", "rb") as image_file:
            #     excel_image_base64 = base64.b64encode(image_file.read()).decode()
            
            # # Help text for the Excel button
            # help_text = "Click here to download the Change_log in Excel"
            
            # # Create the HTML for the Excel download button with image and help text
            # excel_button_html = f"""
            # <a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{excel_data_base64}" download="Change_log.xlsx" class="excel-download-button" title="{help_text}">
            #     <img src="data:image/png;base64,{excel_image_base64}" alt="Download Excel">
            # </a>
            # """
            
            # # Display the custom Excel download button in Streamlit
            # st.markdown(excel_button_html, unsafe_allow_html=True)
            
            ppt_data_overview = ppt_data_change_log = None
        
            ppt_data_change_log = create_download_link_for_change_log_ppt(st.session_state.df_change_log)
            
            # Store ppt_data_change_log in session state if not already stored
            st.session_state.ppt_data_change_log = ppt_data_change_log
                
            if 'ppt_data_overview' in st.session_state:
                ppt_data_overview = st.session_state.ppt_data_overview
            
            if ppt_data_overview and ppt_data_change_log:
                presentation1 = load_presentation_from_bytesio(ppt_data_overview)
                presentation2 = load_presentation_from_bytesio(ppt_data_change_log)
                
                merged_presentation = merge_presentations(presentation1, presentation2)
                merged_presentation_bytesio = BytesIO()
                merged_presentation.save(merged_presentation_bytesio)
                merged_presentation_bytesio.seek(0)
            
                # Base directory of the current script
                base_dir = os.path.dirname(__file__)
                
                # Construct the path for the image
                image_path = os.path.join(base_dir, "Images", "ppt_logo.png")
                
                # Read and encode the image to base64
                with open(image_path, "rb") as image_file:
                    image_base64 = base64.b64encode(image_file.read()).decode()
                        
                # Convert BytesIO to bytes
                ppt_data_merged_bytes = merged_presentation_bytesio.getvalue()
                
                # Encode the PPT data to base64
                ppt_data_merged_base64 = base64.b64encode(ppt_data_merged_bytes).decode()
                
                # Help text
                help_text = "Click here to download the Dashboard into PowerPoint presentation"
                
                # Create the HTML for the button with image and help text
                ppt_button_html = f"""
                <a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{ppt_data_merged_base64}" download="Dashboard.pptx" class="ppt-download-button" title="{help_text}">
                    <img src="data:image/png;base64,{image_base64}" alt="Download PPT">
                </a>
                """
                
                # Display the custom download button in Streamlit
                st.markdown(ppt_button_html, unsafe_allow_html=True)
                
                
                # Check if the updated workbook is available in session state
                if "updated_workbook" in st.session_state:
                    buffer = st.session_state.updated_workbook
                
                    # Load the workbook from the buffer
                    book = load_workbook(buffer)
                
                    # Access the change log sheet and modify cell A26
                    if "1. Change_Log" in book.sheetnames:
                        change_log_sheet = book["1. Change_Log"]
                        
                        # Check if new_entry_df is stored in session state
                        if 'new_entry_df' in st.session_state and st.session_state.entry_added:
                            new_entry_df = st.session_state.new_entry_df
                            
                            # Append new entry to the change log sheet
                            for i, row in new_entry_df.iterrows():
                                change_log_sheet.append(row.tolist())
                
                        # Save the workbook back to the buffer
                        buffer = BytesIO()
                        book.save(buffer)
                        buffer.seek(0)
                
                        # Update the session state with the modified workbook
                        st.session_state.updated_workbook = buffer
                        
                        # Convert BytesIO to bytes
                        excel_data_bytes = buffer.getvalue()
                    
                        # Encode the Excel data to base64
                        excel_data_base64 = base64.b64encode(excel_data_bytes).decode()
                    
                        # Base directory of the current script
                        base_dir = os.path.dirname(__file__)
                        
                        # Construct the path for the image
                        image_path = os.path.join(base_dir, "Images", "excel_logo.png")
                        
                        # Read and encode the image to base64
                        with open(image_path, "rb") as image_file:
                            excel_image_base64 = base64.b64encode(image_file.read()).decode()
                    
                        # Help text for the Excel button
                        help_text = "Click here to download the Consolidated Excel Workbook"
                    
                        # Create the HTML for the Excel download button with image and help text
                        excel_button_html = f"""
                        <a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{excel_data_base64}" download="Excel_template.xlsx" class="excel-download-button" title="{help_text}">
                            <img src="data:image/png;base64,{excel_image_base64}" alt="Download Excel">
                        </a>
                        """
                        #Saving button in session state so that we can use it in another module
                        st.session_state.excel_button_html = excel_button_html
                        
                        custom_css = """
                            <style>
                                .excel-download-button {
                                    position: absolute;
                                    top: -60px;
                                    left: 95px;
                                    cursor: pointer;
                                    z-index: 10;
                                }
                            </style>
                            """
                        st.markdown(custom_css, unsafe_allow_html=True)
                        
                        # Display the custom Excel download button in Streamlit
                        st.markdown(excel_button_html, unsafe_allow_html=True)
                        
                        # Reset the entry_added flag after saving
                        st.session_state.entry_added = False
                     

        else:
            selected_app = next((a for a in self.apps if a["title"] == app), None)
            if selected_app:
                selected_app["function"].app()


if __name__ == "__main__":
    multi_app = MultiApp()
    
    multi_app.add_app("Summary", Summary)
    multi_app.add_app("Gini", gini)
    multi_app.add_app("Calibration", Calibration)
    multi_app.add_app("PSI", PSI)
    multi_app.add_app("Data", Data)
    multi_app.add_app("PPT Customization", Customization)

    multi_app.run()
