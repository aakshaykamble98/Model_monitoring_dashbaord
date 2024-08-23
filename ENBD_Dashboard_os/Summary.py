import streamlit as st
import pandas as pd
import io
import os
import base64
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
import pickle

# from Code.gini import highlight_gini
# from Code.Calibration import highlight_gini_threshold1_calibration, highlight_gini_threshold2_calibration
# from Code.PSI import highlight_gini_PSI

from gini import highlight_gini
from Calibration import highlight_gini_threshold1_calibration, highlight_gini_threshold2_calibration
from PSI import highlight_gini_PSI

# Function to load a PowerPoint presentation from BytesIO
def load_presentation_from_bytesio(presentation_bytesio):
    return Presentation(presentation_bytesio)

# Function to apply title styling
def style_title(title_shape):
    title_shape.left = Inches(0.5)
    title_shape.top = Inches(0.2)
    title_shape.width = Inches(9)
    title_shape.height = Inches(0.5)
    title_text_frame = title_shape.text_frame
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(22)
            # Check if color is in session state; if not, use default color
            if 'title_font_color' in st.session_state:
                run.font.color.rgb = RGBColor.from_string(st.session_state.title_font_color[1:])
            else:
                run.font.color.rgb = RGBColor(0, 0, 0)  # Default color black

# Function to apply title styling and background color
def set_slide_background_and_title_style(slide, title_shape, slide_index):
    
    # Create presentation object
    prs = Presentation()
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    # Check if color is in session state; if not, use default color
    if 'bg_color' in st.session_state:
        fill.fore_color.rgb = RGBColor.from_string(st.session_state.bg_color[1:])
    else:
        fill.fore_color.rgb = RGBColor(0x0A, 0x3D, 0x62)
    
    #Add slide number to title slide
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    ribbon_height = Inches(0.28)  # Adjust height as needed
    half_ribbon_width = slide_width / 2
    ribbon_top = slide_height - ribbon_height
    
    # Add the second half of the ribbon (dark blue)
    ribbon_left = half_ribbon_width  # Adjust the left position for the second half
    shape_blue = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, ribbon_left, ribbon_top, half_ribbon_width, ribbon_height
    )
    fill_blue = shape_blue.fill
    fill_blue.solid()
    # Check if olor is in session state; if not, use default color
    if 'bg_color' in st.session_state:
        fill_blue.fore_color.rgb = RGBColor.from_string(st.session_state.bg_color[1:])  # Dark blue color
    else:
        fill_blue.fore_color.rgb = RGBColor(0x0A, 0x3D, 0x62)  # Dark blue color
    line_blue = shape_blue.line
    # Check if title_font_color is in session state; if not, use default color black
    if 'bg_color' in st.session_state:
        line_blue.color.rgb = RGBColor.from_string(st.session_state.bg_color[1:]) # Dark blue outline color
    else:
        line_blue.color.rgb = RGBColor(0x0A, 0x3D, 0x62) # Dark blue outline color
    
    text_frame = shape_blue.text_frame
    text_frame.text = f" {slide_index}\t"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(9)
    # Check if color is in session state; if not, use default color
    if 'ribbon_font_color' in st.session_state:
        p.font.color.rgb = RGBColor.from_string(st.session_state.ribbon_font_color[1:])  # White text
    else:
        p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    p.alignment = PP_ALIGN.RIGHT  # Align text to the right

    # Apply title styling
    title_shape.left = Inches(1)
    title_shape.top = Inches(2.6)  # Centered vertically
    title_shape.width = Inches(8)
    title_shape.height = Inches(2)
    title_text_frame = title_shape.text_frame
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(42)
            run.font.bold = True
            # Check if color is in session state; if not, use default color
            if 'font_color' in st.session_state:
                run.font.color.rgb = RGBColor.from_string(st.session_state.font_color[1:])  # White font color
            else:
                run.font.color.rgb = RGBColor(255, 255, 255)  # White font color
            
            
# Function to clone shapes (including GraphicFrame and pictures)
def clone_shape(original_shape, merged_slide):
    if original_shape.shape_type == 14:  # Placeholder
        ph = original_shape.placeholder_format
        merged_slide.placeholders[ph.idx].text = original_shape.text
    elif original_shape.shape_type == 13:  # Picture
        image_stream = BytesIO(original_shape.image.blob)
        left = original_shape.left
        top = original_shape.top
        width = original_shape.width
        height = original_shape.height
        merged_slide.shapes.add_picture(image_stream, left, top, width, height)
    else:  # Other shapes (including GraphicFrame)
        merged_shape = original_shape._element
        merged_slide.shapes._spTree.append(merged_shape)


# Function to load thresholds from file
def load_thresholds_gini(file_name='model_gini.pkl'):
    # Base directory of the current script
    base_dir = os.path.dirname(__file__)

    # Construct the path for the pickle file
    file_path = os.path.join(base_dir, 'pkl', file_name)

    try:
        with open(file_path, 'rb') as f:
            thresholds_gini = pickle.load(f)
    except FileNotFoundError:
        thresholds_gini = None

    return thresholds_gini

# Function to load thresholds from file
def load_thresholds_calibration(file_name='model_calibration.pkl'):
    # Base directory of the current script
    base_dir = os.path.dirname(__file__)

    # Construct the path for the pickle file
    file_path = os.path.join(base_dir, 'pkl', file_name)

    try:
        with open(file_path, 'rb') as j:
            thresholds_calibration = pickle.load(j)
    except FileNotFoundError:
        thresholds_calibration = None

    return thresholds_calibration

# Function to load thresholds from file
def load_thresholds_psi(file_name='model_psi.pkl'):
    # Base directory of the current script
    base_dir = os.path.dirname(__file__)

    # Construct the path for the pickle file
    file_path = os.path.join(base_dir, 'pkl', file_name)

    try:
        with open(file_path, 'rb') as k:
            thresholds_psi = pickle.load(k)
    except FileNotFoundError:
        thresholds_psi = None

    return thresholds_psi

# Function to add a slide with the ribbon and logo
def ppt_ribbon_and_logo(slide, slide_index):
    # Create presentation object
    prs = Presentation()
    # Determine slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Calculate ribbon dimensions
    ribbon_height = Inches(0.28)  # Adjust height as needed
    half_ribbon_width = slide_width / 2

    # Calculate positions for ribbons
    ribbon_left = Inches(0)
    ribbon_top = slide_height - ribbon_height

    # Add the first half of the ribbon (amber)
    shape_amber = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, ribbon_left, ribbon_top, half_ribbon_width, ribbon_height
    )
    fill_amber = shape_amber.fill
    fill_amber.solid()
    # Check if color is in session state; if not, use default color
    if 'ribbon_color_1' in st.session_state:
        fill_amber.fore_color.rgb = RGBColor.from_string(st.session_state.ribbon_color_1[1:])  # Amber color
    else:
        fill_amber.fore_color.rgb = RGBColor(255, 191, 0)  # Amber color

    line_amber = shape_amber.line
    # Check if color is in session state; if not, use default color
    if 'ribbon_color_1' in st.session_state:
        line_amber.color.rgb = RGBColor.from_string(st.session_state.ribbon_color_1[1:])  # Amber outline color
    else:
        line_amber.color.rgb = RGBColor(255, 191, 0)  # Amber color

    # Add text to the amber ribbon
    text_frame_amber = shape_amber.text_frame
    text_frame_amber.text = "\tModel Monitoring"
    text_frame_amber.paragraphs[0].font.size = Pt(9)
    # Check if title_font_color is in session state; if not, use default color black
    if 'ribbon_font_1' in st.session_state:
        text_frame_amber.paragraphs[0].font.color.rgb = RGBColor.from_string(st.session_state.ribbon_font_1[1:])  # Black text
    else:
        text_frame_amber.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text

    # Add the second half of the ribbon (dark blue)
    ribbon_left = half_ribbon_width  # Adjust the left position for the second half
    shape_blue = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, ribbon_left, ribbon_top, half_ribbon_width, ribbon_height
    )
    fill_blue = shape_blue.fill
    fill_blue.solid()
    # Check if olor is in session state; if not, use default color
    if 'ribbon_color_2' in st.session_state:
        fill_blue.fore_color.rgb = RGBColor.from_string(st.session_state.ribbon_color_2[1:])  # Dark blue color
    else:
        fill_blue.fore_color.rgb = RGBColor(255, 191, 0)   # Dark blue color
    line_blue = shape_blue.line
    # Check if title_font_color is in session state; if not, use default color black
    if 'ribbon_color_2' in st.session_state:
        line_blue.color.rgb = RGBColor.from_string(st.session_state.ribbon_color_2[1:]) # Dark blue outline color
    else:
        line_blue.color.rgb = RGBColor(255, 191, 0)  # Dark blue outline color
    
    # Add slide number
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    ribbon_height = Inches(0.28)  # Adjust height as needed
    half_ribbon_width = slide_width / 2
    ribbon_top = slide_height - ribbon_height
    
    text_frame = shape_blue.text_frame
    text_frame.text = f" {slide_index}\t"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(9)
    # Check if color is in session state; if not, use default color
    if 'ribbon_font_2' in st.session_state:
        p.font.color.rgb = RGBColor.from_string(st.session_state.ribbon_font_2[1:])  # White text
    else:
        p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    p.alignment = PP_ALIGN.RIGHT  # Align text to the right
    
    # Add logo image to the upper right corner
    base_dir = os.path.dirname(__file__)
    
    # Logo path
    logo_path = os.path.join(base_dir, 'Images', 'ENBD_s.jpg')  # Replace with your actual file path
    logo_left = slide_width - Inches(0.6)  # Adjust position as needed
    logo_top = Inches(0.15)  # Adjust position as needed
    logo_height = Inches(0.45)  # Adjust size as needed
    slide.shapes.add_picture(logo_path, logo_left, logo_top, height=logo_height)

# Creating PowerPoint
def generate_powerpoint_summary(df, thresholds_gini, thresholds_calibration, thresholds_psi):
    prs = Presentation()
    
    slide_index = 1  # To keep track of the slide index
    
    # Add a slide with title and content layout
    slide_layout = prs.slide_layouts[5]  # Use the blank layout
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    ppt_ribbon_and_logo(slide, slide_index)
    
    # Title slide
    title_shape = shapes.title
    title_shape.text = "Summary Table"
    title_shape.left = Inches(0.5)
    title_shape.top = Inches(0.2)
    title_shape.width = Inches(9)
    title_shape.height = Inches(0.5)
    title_text_frame = title_shape.text_frame
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(22)
            # Check if color is in session state; if not, use default color
            if 'title_font_color' in st.session_state:
                run.font.color.rgb = RGBColor.from_string(st.session_state.title_font_color[1:])
            else:
                run.font.color.rgb = RGBColor(0, 0, 0)
    
    # Determine table size and position
    slide_width = prs.slide_width
    table_width = Inches(7.2)  # Adjust width to fit content
    table_height = Inches(5)  # Adjusted to fit the table within the slide
    left = (slide_width - table_width) / 2  # Center the table horizontally
    top = Inches(2.2)  # Center the table vertically

    rows, cols = df.shape
    table = shapes.add_table(rows + 1, cols, left, top, table_width, table_height).table
    table_style_id = table._tbl.tblPr.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}tableStyleId"
    )
    table_style_id.text = "{5940675A-B579-460E-94D1-54222C63F5DA}"
            
    
    # Set column names and font size
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = col_name
        cell.fill.solid()
        # Check if color is in session state; if not, use default color
        if 'row_bg_color' in st.session_state:
            cell.fill.fore_color.rgb = RGBColor.from_string(st.session_state.row_bg_color[1:]) # Teal
        else:
            cell.fill.fore_color.rgb = RGBColor(0x00, 0x80, 0x80) # Teal
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(12)
                # Check if color is in session state; if not, use default color
                if 'row_font_color' in st.session_state:
                    run.font.color.rgb = RGBColor.from_string(st.session_state.row_font_color[1:])  # White font color
                else:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)     # White font color
    
    # Populating the table with DataFrame values and applying highlighting
    for row_idx, row in df.iterrows():
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            if isinstance(value, float):
                cell.text = f"{value:.4f}".rstrip('0').rstrip('.')
            else:
                cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    # Check if color is in session state; if not, use default color
                    if 'content_font_color' in st.session_state:
                        run.font.color.rgb = RGBColor.from_string(st.session_state.content_font_color[1:])
                    else:
                        run.font.color.rgb = RGBColor(0, 0, 0)
            
            if df.columns[col_idx] == 'Dev Gini':
                if value > thresholds_gini['green_gini']['value']:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x00, 0xFF, 0x00)
                elif thresholds_gini['amber_gini']['lower'] < value <= thresholds_gini['amber_gini']['upper']:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xBF, 0x00)
                elif value <= thresholds_gini['red_gini']['value']:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0x00)
            
            elif df.columns[col_idx] == 'Calibration':
                if value < 0:
                    if value > thresholds_calibration['green_calibration_1']['value']:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0x00, 0xFF, 0x00)
                    elif thresholds_calibration['amber_calibration_1']['lower'] <= value <= thresholds_calibration['amber_calibration_1']['upper']:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0xFF, 0xBF, 0x00)
                    elif value < thresholds_calibration['red_calibration_1']['value']:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0x00)
                else:
                    if value < thresholds_calibration['green_calibration_2']['value']:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0x00, 0xFF, 0x00)
                    elif thresholds_calibration['amber_calibration_2']['lower'] <= value <= thresholds_calibration['amber_calibration_2']['upper']:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0xFF, 0xBF, 0x00)
                    elif value > thresholds_calibration['red_calibration_2']['value']:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0x00)
            
            elif df.columns[col_idx] == 'PSI':
                try:
                    value = float(value)
                except ValueError:
                    continue
                if value <= thresholds_psi['green_psi']['value']:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x00, 0xFF, 0x00)
                elif thresholds_psi['amber_psi']['lower'] < value <= thresholds_psi['amber_psi']['upper']:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xBF, 0x00)
                elif value > thresholds_psi['red_psi']['value']:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0x00)
                    
    # Set row heights
    for row in range(rows + 1):
        table.rows[row].height = Inches(0.3)
 
    return prs

# Function to create a download link for the Summary presentation
def create_download_link_for_summary_ppt(df, thresholds_gini, thresholds_calibration, thresholds_psi):
    presentation = generate_powerpoint_summary(df, thresholds_gini, thresholds_calibration, thresholds_psi)
    ppt_stream = BytesIO()
    presentation.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

# Function to merge four PowerPoint presentations and apply title styling to specific slides
def merge_presentations(presentation1, presentation2, presentation3):
    merged_presentation = Presentation()
    slide_index = 1  # To keep track of the slide index

    # Copy slides from presentation1 (Overview) first
    for slide in presentation1.slides:
        slide_layout = merged_presentation.slide_master.slide_layouts[5]
        merged_slide = merged_presentation.slides.add_slide(slide_layout)

        for shape in slide.shapes:
            clone_shape(shape, merged_slide)

        # Apply title styling if there is a title and the slide index is 1
        if merged_slide.shapes.title:
            if slide_index in [1]:
                style_title(merged_slide.shapes.title)
                ppt_ribbon_and_logo(merged_slide, slide_index)
        
        slide_index += 1
        
    # Copy slides from presentation2 (Change Log)
    for slide in presentation2.slides:
        slide_layout = merged_presentation.slide_master.slide_layouts[5]
        merged_slide = merged_presentation.slides.add_slide(slide_layout)

        for shape in slide.shapes:
            clone_shape(shape, merged_slide)

        # Apply title styling if there is a title and the slide index is 1
        if merged_slide.shapes.title:
            if slide_index in [2]:
                style_title(merged_slide.shapes.title)
                ppt_ribbon_and_logo(merged_slide, slide_index)
        
        slide_index += 1
        
    # Copy slides from presentation3 (Summary)
    for slide in presentation3.slides:
        slide_layout = merged_presentation.slide_master.slide_layouts[5]
        merged_slide = merged_presentation.slides.add_slide(slide_layout)

        for shape in slide.shapes:
            clone_shape(shape, merged_slide)

        # Apply title styling if there is a title and the slide index is 3
        if merged_slide.shapes.title:
            if slide_index in [3]:
                style_title(merged_slide.shapes.title)
                ppt_ribbon_and_logo(merged_slide, slide_index)
      
        slide_index += 1
        
    return merged_presentation

# Streamlit app
def app():
    st.markdown(
        """
        <div style='background-color: #008080; border-radius: 5px; padding: 1px;'>
            <h1 style='text-align: center; font-size: 28px; color: white;'>PL - Scorecard Model Summary</h1>
        </div>
        """,
        unsafe_allow_html=True
    )
    
    # Define your CSS styles with custom colors
    css = """
        <style>
            .custom-heading {
                font-weight: bold;
                font-size: 20px; /* Adjust size as needed */
                color: #44546A; /* Custom color for the heading */
                margin-bottom: 5px; /* Adjust spacing between headings */
            }
            .custom-subheading {
                font-weight: bold;
                font-size: 17px; /* Adjust size as needed */
                color: #4472C4; /* Custom color for the subheading */
                margin-bottom: 3px; /* Adjust spacing between subheading and date */
            }
            .custom-date {
                font-weight: bold;
                font-size: 16px; /* Adjust size as needed */
            }
        </style>
    """
    st.markdown(css, unsafe_allow_html=True)
    
    # Initialize variables with default values
    heading = ""
    subheading = ""
    monitoring_date = ""
    date = ""
        
    if 'heading' in st.session_state:
        heading = st.session_state.heading
        
    if 'subheading' in st.session_state:
        subheading = st.session_state.subheading
    
    if 'monitoring_date' in st.session_state:
        monitoring_date = st.session_state.monitoring_date
    
    if 'date' in st.session_state:
        date = st.session_state.date
    
    # Display the extracted heading with custom CSS classes
    st.markdown(f"""
        <div class="custom-heading">{heading}</div>
        <div class="custom-date">{monitoring_date} {date}</div>
    """, unsafe_allow_html=True)

    
    # Load thresholds from file
    thresholds_gini = load_thresholds_gini()
    if thresholds_gini is None:
        st.error("Thresholds not found! Please set thresholds in the Gini module first.")
        st.stop()
        
    # Load thresholds from file
    thresholds_calibration = load_thresholds_calibration()
    if thresholds_calibration is None:
        st.error("Thresholds not found! Please set thresholds in the calibration module first.")
        st.stop()
        
    # Load thresholds from file
    thresholds_psi = load_thresholds_psi()
    if thresholds_psi is None:
        st.error("Thresholds not found! Please set thresholds in the PSI module first.")
        st.stop()
    
    
    # Base directory of the current script
    base_dir = os.path.dirname(__file__)
    
    # Read the Excel file
    df = pd.read_excel(os.path.join(base_dir, 'Datasets', 'Summary_table.xlsx'))

    # Apply styling to the DataFrame
    styled_df = df.style.applymap(lambda x: highlight_gini(x, thresholds_gini), subset=['Dev Gini'])

    # Apply conditional highlighting based on the 'Calibration' column
    styled_df = styled_df.applymap(lambda x: highlight_gini_threshold1_calibration(x, thresholds_calibration) if x < 0 else highlight_gini_threshold2_calibration(x, thresholds_calibration) if x > 0 else x, subset=['Calibration'])
    
    # Highlight PSI row in PD Bucket column
    styled_df = styled_df.applymap(lambda val: highlight_gini_PSI(val, thresholds_psi), subset=['PSI'])
    
    # Display the DataFrame without unnecessary trailing zeros
    df_styled = styled_df.format(lambda x: f"{x:.4f}".rstrip('0').rstrip('.') if isinstance(x, float) else f"{x}")
    st.dataframe(df_styled, width=1200)

# =============================================================================
#     # Prepare Excel file in memory
#     excel_data = io.BytesIO()
#     with pd.ExcelWriter(excel_data, engine='openpyxl') as writer:
#         df_styled.to_excel(writer, index=False, sheet_name='Summary', startrow=12)
#     excel_data.seek(0)
#     
#     # Save byte data into st.session_state
#     st.session_state.summary_workbook_data = excel_data
# 
#     # Convert BytesIO to bytes
#     excel_data_bytes = excel_data.getvalue()
#     
#     # Encode the Excel data to base64
#     excel_data_base64 = base64.b64encode(excel_data_bytes).decode()
#     
#     # Read and encode the Excel image to base64
#     with open("excel_logo.png", "rb") as image_file:
#         excel_image_base64 = base64.b64encode(image_file.read()).decode()
#     
#     # Help text for the Excel button
#     help_text = "Click here to download the summary in Excel"
#     
#     # Create the HTML for the Excel download button with image and help text
#     excel_button_html = f"""
#     <a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{excel_data_base64}" download="Summary_table.xlsx" class="excel-download-button" title="{help_text}">
#         <img src="data:image/png;base64,{excel_image_base64}" alt="Download Excel">
#     </a>
#     """
#     
#     # Display the custom Excel download button in Streamlit
#     st.markdown(excel_button_html, unsafe_allow_html=True)
# =============================================================================

    ppt_bytes = create_download_link_for_summary_ppt(df, thresholds_gini, thresholds_calibration, thresholds_psi)
    
    # Save ppt_bytes in session state
    st.session_state.ppt_data_summary = ppt_bytes
    
    ppt_data_overview = ppt_data_change_log = ppt_data_summary = None
    
    if 'ppt_data_overview' in st.session_state:
        ppt_data_overview = st.session_state.ppt_data_overview
    
    if 'ppt_data_change_log' in st.session_state:
        ppt_data_change_log = st.session_state.ppt_data_change_log
    
    if ppt_data_overview and ppt_data_change_log and ppt_bytes:
        presentation1 = load_presentation_from_bytesio(ppt_data_overview)
        presentation2 = load_presentation_from_bytesio(ppt_data_change_log)
        presentation3 = load_presentation_from_bytesio(ppt_bytes)
        
        merged_presentation = merge_presentations(presentation1, presentation2, presentation3)
        merged_presentation_bytesio = BytesIO()
        merged_presentation.save(merged_presentation_bytesio)
        merged_presentation_bytesio.seek(0)
    
        # Base directory of the current script
        base_dir = os.path.dirname(__file__)
        
        # Construct the path for the image
        image_path = os.path.join(base_dir, "Images", "ppt_logo.png")
        
        # Read and encode the image to base64
        with open(image_path, "rb") as image_file:
            image_base64 = base64.b64encode(image_file.read()).decode()
        
        # Convert BytesIO to bytes
        ppt_data_merged_bytes = merged_presentation_bytesio.getvalue()
        
        # Encode the PPT data to base64
        ppt_data_merged_base64 = base64.b64encode(ppt_data_merged_bytes).decode()
        
        # Help text
        help_text = "Click here to download the Dashboard into PowerPoint presentation"
        
        # Create the HTML for the button with image and help text
        ppt_button_html = f"""
        <a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{ppt_data_merged_base64}" download="Dashboard.pptx" class="ppt-download-button" title="{help_text}">
            <img src="data:image/png;base64,{image_base64}" alt="Download PPT">
        </a>
        """
        
        # Display the custom download button in Streamlit
        st.markdown(ppt_button_html, unsafe_allow_html=True)
        
        custom_css = """
            <style>
                .excel-download-button {
                    position: absolute;
                    top: -60px;
                    left: 95px;
                    cursor: pointer;
                    z-index: 10;
                }
            </style>
            """
        st.markdown(custom_css, unsafe_allow_html=True)
        
    # Ensure that you check if the button HTML is available in session state
    if 'excel_button_html' in st.session_state:
        st.markdown(st.session_state.excel_button_html, unsafe_allow_html=True)
    else:
        st.write("No data available for download in Excel. Please run the Overview and Change Log modules first.")

if __name__ == "__main__":
    app()
