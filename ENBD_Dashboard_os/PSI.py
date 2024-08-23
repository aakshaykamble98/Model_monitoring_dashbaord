import streamlit as st
import pandas as pd
import numpy as np
import plotly.graph_objects as go
import plotly.io as pio
import io
import os
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from openpyxl import Workbook
from openpyxl.styles import PatternFill
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from streamlit_modal import Modal
from io import BytesIO
import base64
import pickle
from gini import create_ppt_download_button_gini
from Calibration import create_ppt_download_button_calibration


# Function to load a PowerPoint presentation from BytesIO
def load_presentation_from_bytesio(presentation_bytesio):
    return Presentation(presentation_bytesio)

# Function to apply title styling
def style_title(title_shape):
    title_shape.left = Inches(0.5)
    title_shape.top = Inches(0.2)
    title_shape.width = Inches(9)
    title_shape.height = Inches(0.5)
    title_text_frame = title_shape.text_frame
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(22)
            # Check if color is in session state; if not, use default color
            if 'title_font_color' in st.session_state:
                run.font.color.rgb = RGBColor.from_string(st.session_state.title_font_color[1:])
            else:
                run.font.color.rgb = RGBColor(0, 0, 0)  # Default color black

# Function to apply title styling and background color
def set_slide_background_and_title_style(slide, title_shape, slide_index):
    
    # Create presentation object
    prs = Presentation()
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    # Check if color is in session state; if not, use default color
    if 'bg_color' in st.session_state:
        fill.fore_color.rgb = RGBColor.from_string(st.session_state.bg_color[1:])
    else:
        fill.fore_color.rgb = RGBColor(0x00, 0x80, 0x80)
    
    #Add slide number to title slide
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    ribbon_height = Inches(0.28)  # Adjust height as needed
    ribbon_top = slide_height - ribbon_height
    
    # Add the second half of the ribbon (dark blue)
    ribbon_left = 0  # Adjust the position 
    shape_blue = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, ribbon_left, ribbon_top, slide_width, ribbon_height
    )
    fill_blue = shape_blue.fill
    fill_blue.solid()
    # Check if olor is in session state; if not, use default color
    if 'bg_color' in st.session_state:
        fill_blue.fore_color.rgb = RGBColor.from_string(st.session_state.bg_color[1:])  # Dark blue color
    else:
        fill_blue.fore_color.rgb = RGBColor(0x00, 0x80, 0x80)  # Dark blue color
    line_blue = shape_blue.line
    # Check if title_font_color is in session state; if not, use default color black
    if 'bg_color' in st.session_state:
        line_blue.color.rgb = RGBColor.from_string(st.session_state.bg_color[1:]) # Dark blue outline color
    else:
        line_blue.color.rgb = RGBColor(0x00, 0x80, 0x80) # Dark blue outline color
    

    text_frame = shape_blue.text_frame
    text_frame.text = f" {slide_index}\t"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(9)
    # Check if color is in session state; if not, use default color
    if 'ribbon_font_color' in st.session_state:
        p.font.color.rgb = RGBColor.from_string(st.session_state.ribbon_font_color[1:])  # White text
    else:
        p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    p.alignment = PP_ALIGN.RIGHT  # Align text to the right

    # Apply title styling
    title_shape.left = Inches(1)
    title_shape.top = Inches(2.6)  # Centered vertically
    title_shape.width = Inches(8)
    title_shape.height = Inches(2)
    title_text_frame = title_shape.text_frame
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(42)
            run.font.bold = True
            # Check if color is in session state; if not, use default color
            if 'font_color' in st.session_state:
                run.font.color.rgb = RGBColor.from_string(st.session_state.font_color[1:])  # White font color
            else:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White font color
            
            
# Function to clone shapes (including GraphicFrame and pictures)
def clone_shape(original_shape, merged_slide):
    if original_shape.shape_type == 14:  # Placeholder
        ph = original_shape.placeholder_format
        merged_slide.placeholders[ph.idx].text = original_shape.text
    elif original_shape.shape_type == 13:  # Picture
        image_stream = BytesIO(original_shape.image.blob)
        left = original_shape.left
        top = original_shape.top
        width = original_shape.width
        height = original_shape.height
        merged_slide.shapes.add_picture(image_stream, left, top, width, height)
    else:  # Other shapes (including GraphicFrame)
        merged_shape = original_shape._element
        merged_slide.shapes._spTree.append(merged_shape)

# Function to highlight Gini values based on user-defined thresholds
def highlight_gini_PSI(val, thresholds_psi):
    try:
        if val == '' or pd.isna(val):
            color = 'white'
        else:
            val = float(val)
            if val <= thresholds_psi['green_psi']['value']:
                color = 'green'
            elif thresholds_psi['amber_psi']['lower'] < val <= thresholds_psi['amber_psi']['upper']:
                color = 'orange'
            elif val > thresholds_psi['red_psi']['value']:
                color = 'red'
            else:
                color = 'white'
    except ValueError:
        color = 'white'
    return f'background-color: {color}'

# # Function to create an Excel file with highlighted cells
# def to_excel_with_highlights_PSI(df, thresholds_psi):
#     output = io.BytesIO()
#     workbook = Workbook()
#     sheet = workbook.active

#     # Write the header
#     for col_num, column_title in enumerate(df.columns, 1):
#         cell = sheet.cell(row=1, column=col_num)
#         cell.value = column_title

#     # Write the data
#     for row_num, row in enumerate(df.values, 2):
#         for col_num, value in enumerate(row, 1):
#             cell = sheet.cell(row=row_num, column=col_num)
#             cell.value = value

#     # Highlight the row named 'PSI' in the 'PD Bucket' column
#     for row in sheet.iter_rows(min_row=2, max_row=sheet.max_row):
#         if row[df.columns.get_loc('PD Bucket')].value == 'PSI':
#             for cell in row:
#                 try:
#                     value = float(cell.value)
#                 except ValueError:
#                     continue  # Skip cells that can't be converted to float
#                 if value <= thresholds_psi['green_psi']['value']:
#                     cell.fill = PatternFill(start_color="00FF00", end_color="00FF00", fill_type="solid")
#                 elif thresholds_psi['amber_psi']['lower'] < value <= thresholds_psi['amber_psi']['upper']:
#                     cell.fill = PatternFill(start_color="FFBF00", end_color="FFBF00", fill_type="solid")
#                 elif value > thresholds_psi['red_psi']['value']:
#                     cell.fill = PatternFill(start_color="FF0000", end_color="FF0000", fill_type="solid")

#     workbook.save(output)
#     #st.session_state.psi_workbook_data = output.getvalue()
#     return output.getvalue()

# Function to add a slide with the ribbon and logo
def ppt_ribbon_and_logo(slide, slide_index):
    # Create presentation object
    prs = Presentation()
    # Determine slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Calculate ribbon dimensions
    ribbon_height = Inches(0.28)  # Adjust height as needed
    half_ribbon_width = slide_width / 2

    # Calculate positions for ribbons
    ribbon_left = Inches(0)
    ribbon_top = slide_height - ribbon_height

    # Add the first half of the ribbon (amber)
    shape_amber = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, ribbon_left, ribbon_top, half_ribbon_width, ribbon_height
    )
    fill_amber = shape_amber.fill
    fill_amber.solid()
    # Check if color is in session state; if not, use default color
    if 'ribbon_color_1' in st.session_state:
        fill_amber.fore_color.rgb = RGBColor.from_string(st.session_state.ribbon_color_1[1:])  # Amber color
    else:
        fill_amber.fore_color.rgb = RGBColor(255, 191, 0)  # Amber color

    line_amber = shape_amber.line
    # Check if color is in session state; if not, use default color
    if 'ribbon_color_1' in st.session_state:
        line_amber.color.rgb = RGBColor.from_string(st.session_state.ribbon_color_1[1:])  # Amber outline color
    else:
        line_amber.color.rgb = RGBColor(255, 191, 0)  # Amber color

    # Add text to the amber ribbon
    text_frame_amber = shape_amber.text_frame
    text_frame_amber.text = "\tModel Monitoring"
    text_frame_amber.paragraphs[0].font.size = Pt(9)
    # Check if color is in session state; if not, use default color
    if 'ribbon_font_1' in st.session_state:
        text_frame_amber.paragraphs[0].font.color.rgb = RGBColor.from_string(st.session_state.ribbon_font_1[1:])  # Black text
    else:
        text_frame_amber.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text

    # Add the second half of the ribbon (dark blue)
    ribbon_left = half_ribbon_width  # Adjust the left position for the second half
    shape_blue = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, ribbon_left, ribbon_top, half_ribbon_width, ribbon_height
    )
    fill_blue = shape_blue.fill
    fill_blue.solid()
    # Check if color is in session state; if not, use default color
    if 'ribbon_color_2' in st.session_state:
        fill_blue.fore_color.rgb = RGBColor.from_string(st.session_state.ribbon_color_2[1:])  # Dark blue color
    else:
        fill_blue.fore_color.rgb = RGBColor(255, 191, 0)  # Dark blue color

    line_blue = shape_blue.line
    # Check if title_font_color is in session state; if not, use default color black
    if 'ribbon_color_2' in st.session_state:
        line_blue.color.rgb = RGBColor.from_string(st.session_state.ribbon_color_2[1:]) # Dark blue outline color
    else:
        line_blue.color.rgb = RGBColor(255, 191, 0) # Dark blue outline color
    
    # Add slide number
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    ribbon_height = Inches(0.28)  # Adjust height as needed
    half_ribbon_width = slide_width / 2
    ribbon_top = slide_height - ribbon_height
    
    slide_number_box = slide.shapes.add_textbox(half_ribbon_width, ribbon_top, half_ribbon_width, ribbon_height)
    text_frame = slide_number_box.text_frame
    text_frame.text = f" {slide_index}\t"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(9)
    # Check if color is in session state; if not, use default color
    if 'ribbon_font_2' in st.session_state:
        p.font.color.rgb = RGBColor.from_string(st.session_state.ribbon_font_2[1:])  # White text
    else:
        p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    p.alignment = PP_ALIGN.RIGHT  # Align text to the right
    

    # Add logo image to the upper right corner
    base_dir = os.path.dirname(__file__)
    
    # Logo path
    logo_path = os.path.join(base_dir, 'Images', 'ENBD_s.jpg')  # Replace with your actual file path
    logo_left = slide_width - Inches(0.6)  # Adjust position as needed
    logo_top = Inches(0.15)  # Adjust position as needed
    logo_height = Inches(0.45)  # Adjust size as needed
    slide.shapes.add_picture(logo_path, logo_left, logo_top, height=logo_height)


def create_ppt_PSI(df, fig1_bytes, fig2_bytes, thresholds_psi, data_comment, graph_comment):
    prs = Presentation() 
    
    slide_index = 1  # To keep track of the slide index
    
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Set the background color to the slide
    background = slide.background
    fill = background.fill
    fill.solid()
    # Check if color is in session state; if not, use default color
    if 'bg_color' in st.session_state:
        fill.fore_color.rgb = RGBColor.from_string(st.session_state.bg_color[1:])
    else:
        fill.fore_color.rgb = RGBColor(0x00, 0x80, 0x80)
    shapes = slide.shapes
    title_shape = shapes.title
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "PL - Scorecard Model PSI"
    title_shape.left = Inches(1)
    title_shape.top = Inches(2.6)  # Centered vertically
    title_shape.width = Inches(8)
    title_shape.height = Inches(2)
    title_text_frame = title_shape.text_frame
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(42)
            run.font.bold = True
            # Check if color is in session state; if not, use default color
            if 'font_color' in st.session_state:
                run.font.color.rgb = RGBColor.from_string(st.session_state.font_color[1:])  # White font color
            else:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White font color
            
    # Add logo image to the upper right corner
    base_dir = os.path.dirname(__file__)
    
    # Logo path
    logo_path = os.path.join(base_dir, 'Images', 'ENBD.jpg')  # Replace with your actual file path
    logo_left = Inches(0.6)  # Adjust position as needed
    logo_top = Inches(0.25)  # Adjust position as needed
    logo_height = Inches(0.6)  # Adjust size as needed
    slide.shapes.add_picture(logo_path, logo_left, logo_top, height=logo_height)
            
    # Add slide number to title slide        
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    ribbon_height = Inches(0.28)  # Adjust height as needed
    ribbon_top = slide_height - ribbon_height
    
    # Add the second half of the ribbon (dark blue)
    ribbon_left = 0  # Adjust the position
    shape_blue = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, ribbon_left, ribbon_top, slide_width, ribbon_height
    )
    fill_blue = shape_blue.fill
    fill_blue.solid()
    # Check if color is in session state; if not, use default color
    if 'bg_color' in st.session_state:
        fill_blue.fore_color.rgb = RGBColor.from_string(st.session_state.bg_color[1:])  # Dark blue color
    else:
        fill_blue.fore_color.rgb = RGBColor(0x00, 0x80, 0x80)  # Dark blue color
    line_blue = shape_blue.line
    # Check if color is in session state; if not, use default color
    if 'bg_color' in st.session_state:
        line_blue.color.rgb = RGBColor.from_string(st.session_state.bg_color[1:]) # Dark blue outline color
    else:
        line_blue.color.rgb = RGBColor(0x00, 0x80, 0x80) # Dark blue outline color
    
    text_frame = shape_blue.text_frame
    text_frame.text = f" {slide_index}\t"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(9)
    # Check if color is in session state; if not, use default color
    if 'ribbon_font_color' in st.session_state:
        p.font.color.rgb = RGBColor.from_string(st.session_state.ribbon_font_color[1:])
    else:
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White text
    p.alignment = PP_ALIGN.RIGHT  # Align text to the right
    
    slide_index += 1

    # Data Table slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    shapes = slide.shapes
    
    ppt_ribbon_and_logo(slide, slide_index)

    # Add and style title
    title_shape = shapes.title
    title_shape.text = "PSI calculation"
    title_shape.left = Inches(0.5)
    title_shape.top = Inches(0.2)
    title_shape.width = Inches(9)
    title_shape.height = Inches(0.5)
    title_text_frame = title_shape.text_frame
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(22)
            # Check if color is in session state; if not, use default color
            if 'title_font_color' in st.session_state:
                run.font.color.rgb = RGBColor.from_string(st.session_state.title_font_color[1:])
            else:
                run.font.color.rgb = RGBColor(0, 0, 0)

    # Determine table size and position
    left = Inches(0.5)
    top = Inches(0.7)
    width = Inches(9)
    height = Inches(5.0)  # Adjusted to fit the table within the slide

    rows, cols = df.shape
    table = shapes.add_table(rows + 1, cols, left, top, width, height).table
    table_style_id = table._tbl.tblPr.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}tableStyleId"
    )
    table_style_id.text = "{5940675A-B579-460E-94D1-54222C63F5DA}"

    # Set fixed row height
    fixed_row_height = Inches(0)  # Adjust the height as needed

    # Set column names and font size
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = col_name
        cell.fill.solid()
        # Check if color is in session state; if not, use default color
        if 'row_bg_color' in st.session_state:
            cell.fill.fore_color.rgb = RGBColor.from_string(st.session_state.row_bg_color[1:]) # Teal
        else:
            cell.fill.fore_color.rgb = RGBColor(0x00, 0x80, 0x80) # Teal
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(8)
                # Check if color is in session state; if not, use default color
                if 'row_font_color' in st.session_state:
                    run.font.color.rgb = RGBColor.from_string(st.session_state.row_font_color[1:])  # White font color
                else:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)   #White font
                
                
        table.rows[0].height = fixed_row_height  # Set height for the header row

    # Add data to table and set font size
    df.fillna('-', inplace=True)
    for row_idx, row in df.iterrows():
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            # Format decimal values to 2 decimal places if applicable
            if isinstance(value, float):
                cell.text = f"{value:.4f}".rstrip('0').rstrip('.')
            else:
                cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(8)  # Set font size to 8
                    # Check if color is in session state; if not, use default color
                    if 'content_font_color' in st.session_state:
                        run.font.color.rgb = RGBColor.from_string(st.session_state.content_font_color[1:])
                    else:
                        run.font.color.rgb = RGBColor(0, 0, 0)

        table.rows[row_idx + 1].height = fixed_row_height  # Set height for all data rows

    # Highlight PSI row in PD Bucket column
    for row_idx, row in df.iterrows():
        if df.loc[row_idx, 'PD Bucket'] == 'PSI':
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx + 1, col_idx)
                try:
                    psi_value = float(value)
                except ValueError:
                    continue  # Skip cells that can't be converted to float
                if psi_value <= thresholds_psi['green_psi']['value']:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x00, 0xFF, 0x00)
                elif thresholds_psi['amber_psi']['lower'] < psi_value <= thresholds_psi['amber_psi']['upper']:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xBF, 0x00)
                elif psi_value > thresholds_psi['red_psi']['value']:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0x00)

    # Add data comment
    data_comment_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(1))
    text_frame = data_comment_box.text_frame
    text_frame.text = f"Comment: {data_comment}"
    # Set font size to 10 for all paragraphs in the text frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(14)
    
    slide_index += 1

    # First graph slide (Subdivided Bar Plot of PD Buckets)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    ppt_ribbon_and_logo(slide, slide_index)
    
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Graph"
    title_shape.left = Inches(0.5)
    title_shape.top = Inches(0.2)
    title_shape.width = Inches(9)
    title_shape.height = Inches(0.5)
    title_text_frame = title_shape.text_frame
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(22)
            # Check if color is in session state; if not, use default color
            if 'title_font_color' in st.session_state:
                run.font.color.rgb = RGBColor.from_string(st.session_state.title_font_color[1:])
            else:
                run.font.color.rgb = RGBColor(0, 0, 0)

    # Add chart image 1
    image_stream1 = io.BytesIO(fig1_bytes)
    # shapes.add_picture(image_stream1, Inches(1), Inches(0.9), Inches(7.5), Inches(4.5))
    shapes.add_picture(image_stream1, Inches(0.6), Inches(0.8), Inches(8.8), Inches(4.5))
    
    slide_index += 1

    # Second graph slide (Comparison of Dev and 202403)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    ppt_ribbon_and_logo(slide, slide_index)
    
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Graph"
    title_shape.left = Inches(0.5)
    title_shape.top = Inches(0.2)
    title_shape.width = Inches(9)
    title_shape.height = Inches(0.5)
    title_text_frame = title_shape.text_frame
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(22)
            # Check if color is in session state; if not, use default color
            if 'title_font_color' in st.session_state:
                run.font.color.rgb = RGBColor.from_string(st.session_state.title_font_color[1:])
            else:
                run.font.color.rgb = RGBColor(0, 0, 0)

    # Add chart image 2
    image_stream2 = io.BytesIO(fig2_bytes)
    # shapes.add_picture(image_stream2, Inches(1), Inches(0.9), Inches(7.5), Inches(4.5))
    shapes.add_picture(image_stream2, Inches(0.6), Inches(0.8), Inches(8.8), Inches(4.5))

    # Add graph comment
    graph_comment_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    text_frame = graph_comment_box.text_frame
    text_frame.text = f"Comment: {graph_comment}"
    # Set font size to 10 for all paragraphs in the text frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(14)

    # Save presentation
    ppt_output = io.BytesIO()
    prs.save(ppt_output)
    ppt_output.seek(0)
    return ppt_output

# Function to create the threshold selection UI for PSI
def threshold_selection_PSI(show_ui=True):
    # Base directory of the current script
    base_dir = os.path.dirname(__file__)
    
    # Construct the path for the thresholds file
    file_path = os.path.join(base_dir, 'pkl', 'model_psi.pkl')

    try:
        with open(file_path, 'rb') as k:
            thresholds_psi = pickle.load(k)
    except FileNotFoundError:
        thresholds_psi = None
        
    if thresholds_psi:
        green_threshold = thresholds_psi['green_psi']['value']
        amber_lower = thresholds_psi['amber_psi']['lower']
        amber_upper = thresholds_psi['amber_psi']['upper']
        red_threshold = thresholds_psi['red_psi']['value']
    else:
        # Default values when UI is not shown
        green_threshold = None
        amber_lower =None
        amber_upper = None
        red_threshold = None
        
    if show_ui:
        with st.expander('Please select the PSI threshold values'):
            st.markdown('<hr style="margin-top: -5px; margin-bottom: -5px;">', unsafe_allow_html=True)
            c1, c2, c3 = st.columns([2, 4, 2])
            with c1:
                st.markdown('<p style="font-size:17px;"><b>Legend</b></p>', unsafe_allow_html=True)
            with c2:
                st.markdown('<p style="font-size:17px;"><b>Thresholds</b></p>', unsafe_allow_html=True)
            with c3:
                st.markdown('<p style="font-size:17px;"><b>Evaluation</b></p>', unsafe_allow_html=True)
            st.markdown('<hr style="margin-top: -5px; margin-bottom: -5px;">', unsafe_allow_html=True)

            c1, c2, c3 = st.columns([2, 4, 2])
            with c1:
                st.markdown('<div style="background-color: green; width: 80%; height: 40px;"></div>', unsafe_allow_html=True)
            with c2:
                green_threshold = st.number_input("Green Threshold (PSI <= 0.1)", value= green_threshold or 0.1)
            with c3:
                st.markdown('<p style="font-size:19px;"><b>No Action Required</b></p>', unsafe_allow_html=True)
            st.markdown('<hr style="margin-top: -5px; margin-bottom: -5px;">', unsafe_allow_html=True)

            c1, c2, c3, c4 = st.columns(4)
            with c1:
                st.markdown('<div style="background-color: orange; width: 80%; height: 40px;"></div>', unsafe_allow_html=True)
            with c2:
                amber_lower = st.number_input("Amber Lower Limit (0.1 < PSI <= 0.25)", value=amber_lower or 0.1)
            with c3:
                amber_upper = st.number_input("Amber Upper Limit (0.1 < PSI <= 0.25)", value= amber_upper or 0.25 )
            with c4:
                st.markdown('<p style="font-size:19px;"><b>To be Discussed</b></p>', unsafe_allow_html=True)
            st.markdown('<hr style="margin-top: -5px; margin-bottom: -5px;">', unsafe_allow_html=True)

            c1, c2, c3 = st.columns([2, 4, 2])
            with c1:
                st.markdown('<div style="background-color: red; width: 80%; height: 40px;"></div>', unsafe_allow_html=True)
            with c2:
                red_threshold = st.number_input("Red Threshold (PSI > 0.25)", value=red_threshold or 0.25)
            with c3:
                st.markdown('<p style="font-size:19px;"><b>Action Required</b></p>', unsafe_allow_html=True)
            
        # Save thresholds to file
        thresholds_psi = {
            'green_psi': {'value': green_threshold},
            'amber_psi': {'lower': amber_lower, 'upper': amber_upper},
            'red_psi': {'value': red_threshold}
        }
        
        # Base directory of the current script
        base_dir = os.path.dirname(__file__)
        
        # Construct the path for the thresholds file
        file_path = os.path.join(base_dir, 'pkl', 'model_psi.pkl')
    
        with open(file_path, 'wb') as k:
            pickle.dump(thresholds_psi, k)

    return {
        'green_psi': {'value': green_threshold},
        'amber_psi': {'lower': amber_lower, 'upper': amber_upper},
        'red_psi': {'value': red_threshold}
    }


def create_powerpoint_download_button_PSI(df, fig1_bytes, fig2_bytes, thresholds_psi, data_comment="", graph_comment=""):
    # Create PowerPoint presentation bytes
    ppt_data_psi = create_ppt_PSI(df, fig1_bytes, fig2_bytes, thresholds_psi, data_comment, graph_comment)
    
    return ppt_data_psi

# Function to merge four PowerPoint presentations and apply title styling to specific slides
def merge_presentations(presentation1, presentation2, presentation3, presentation4, presentation5, presentation6):
    merged_presentation = Presentation()
    slide_index = 1  # To keep track of the slide index

    # Copy slides from presentation1 (Overview) first
    for slide in presentation1.slides:
        slide_layout = merged_presentation.slide_master.slide_layouts[5]
        merged_slide = merged_presentation.slides.add_slide(slide_layout)

        for shape in slide.shapes:
            clone_shape(shape, merged_slide)

        # Apply title styling if there is a title and the slide index is 1
        if merged_slide.shapes.title:
            if slide_index in [1]:
                style_title(merged_slide.shapes.title)
                ppt_ribbon_and_logo(merged_slide, slide_index)
        
        slide_index += 1
        
    # Copy slides from presentation2 (Change Log)
    for slide in presentation2.slides:
        slide_layout = merged_presentation.slide_master.slide_layouts[5]
        merged_slide = merged_presentation.slides.add_slide(slide_layout)

        for shape in slide.shapes:
            clone_shape(shape, merged_slide)

        # Apply title styling if there is a title and the slide index is 1
        if merged_slide.shapes.title:
            if slide_index in [2]:
                style_title(merged_slide.shapes.title)
                ppt_ribbon_and_logo(merged_slide, slide_index)
        
        slide_index += 1
        
    # Copy slides from presentation3 (Summary)
    for slide in presentation3.slides:
        slide_layout = merged_presentation.slide_master.slide_layouts[5]
        merged_slide = merged_presentation.slides.add_slide(slide_layout)

        for shape in slide.shapes:
            clone_shape(shape, merged_slide)

        # Apply title styling if there is a title and the slide index is 3
        if merged_slide.shapes.title:
            if slide_index in [3]:
                style_title(merged_slide.shapes.title)
                ppt_ribbon_and_logo(merged_slide, slide_index)
      
        slide_index += 1
    
    # Copy slides from presentation4 (Gini)
    for slide in presentation4.slides:
        slide_layout = merged_presentation.slide_master.slide_layouts[5]
        merged_slide = merged_presentation.slides.add_slide(slide_layout)

        for shape in slide.shapes:
            clone_shape(shape, merged_slide)

        # Apply title styling if there is a title and the slide index is 5 or 6
        if merged_slide.shapes.title:
            if slide_index in [5, 6]:
                style_title(merged_slide.shapes.title)
                ppt_ribbon_and_logo(merged_slide, slide_index)
        
            if slide_index in [4]:
                set_slide_background_and_title_style(merged_slide, merged_slide.shapes.title, slide_index)
    
        slide_index += 1

    # Copy slides from presentation5 (Calibration)
    for slide in presentation5.slides:
        slide_layout = merged_presentation.slide_master.slide_layouts[5]
        merged_slide = merged_presentation.slides.add_slide(slide_layout)

        for shape in slide.shapes:
            clone_shape(shape, merged_slide)

        # Apply title styling if there is a title and the slide index is 8 or 9
        if merged_slide.shapes.title:
            if slide_index in [8, 9]:
                style_title(merged_slide.shapes.title)
                ppt_ribbon_and_logo(merged_slide, slide_index)
                
            if slide_index in [7]:
                set_slide_background_and_title_style(merged_slide, merged_slide.shapes.title, slide_index)
      
        slide_index += 1

    # Copy slides from presentation6 (PSI)
    for slide in presentation6.slides:
        slide_layout = merged_presentation.slide_master.slide_layouts[5]
        merged_slide = merged_presentation.slides.add_slide(slide_layout)

        for shape in slide.shapes:
            clone_shape(shape, merged_slide)

        # Apply title styling if there is a title and the slide index is 11 or 12
        if merged_slide.shapes.title:
            if slide_index in [11, 12, 13]:
                style_title(merged_slide.shapes.title)
                ppt_ribbon_and_logo(merged_slide, slide_index)
                
            if slide_index in [10]:
                set_slide_background_and_title_style(merged_slide, merged_slide.shapes.title, slide_index)
       
        slide_index += 1

    return merged_presentation

def app():
    st.markdown(
        """
        <div style='background-color: #008080; border-radius: 5px; padding: 1px;'>
            <h1 style='text-align: center; font-size: 28px; color: white;'>PL - Scorecard Model PSI</h1>
        </div>
        """,
        unsafe_allow_html=True
    )
    
    custom_css = """
    <style>
        .excel-download-button {
            position: absolute;
            top: -60px;
            left: 85px;
            cursor: pointer
        }
    </style>
    """
    st.markdown(custom_css, unsafe_allow_html=True)

    # Base directory of the current script
    base_dir = os.path.dirname(__file__)
    
    # Read the Excel file
    df = pd.read_excel(os.path.join(base_dir, 'Datasets', 'PSI_Data_dashboard.xlsx'))  # Ensure the file is in the 'Datasets' directory
    # Replace None values with empty strings for better display
    df1 = df.fillna("")
    df.columns = [str(i) for i in df.columns]
    st.session_state.df_psi = df  # Save df to session_state

    data_comment_psi = ""
    graph_comment_psi = ""

    tab1, tab2 = st.tabs(["PSI Calculation", "Graph"])

    with tab1:
        thresholds_psi = threshold_selection_PSI(show_ui=True)
        st.session_state.thresholds_psi = thresholds_psi
        
        st.markdown(
                """
                <div style='text-align: center;
                font-size: 20px;'>
                    <strong>PSI Result</strong>
                </div>
                """,
                unsafe_allow_html=True)
        
        # def create_excel_download_button_PSI(df, thresholds_psi):
        #     excel_data = to_excel_with_highlights_PSI(df, thresholds_psi)
            
        #     # Save byte data into st.session_state
        #     st.session_state.psi_workbook_data = excel_data
            
        #     buffer = BytesIO()
        #     buffer.write(excel_data)
        #     buffer.seek(0)
        #     b64 = base64.b64encode(buffer.read()).decode()
            
        #     # Read and encode the Excel image to base64
        #     with open("excel_logo.png", "rb") as image_file:
        #         excel_image_base64 = base64.b64encode(image_file.read()).decode()
            
        #     # Define the HTML for the download icon link
        #     download_html = f'''
        #         <a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64}" 
        #            download="highlighted_data.xlsx" class="excel-download-button" title="Click here to download the Excel file">
        #            <img src="data:image/png;base64,{excel_image_base64}" alt="Download Excel">
        #         </a>
        #     '''
            
        #     return download_html
        
        # download_html = create_excel_download_button_PSI(df, thresholds_psi)

        
        # st.markdown(download_html, unsafe_allow_html=True)

        styled_df = df1.style.applymap(
            lambda val: highlight_gini_PSI(val, thresholds_psi),
            subset=pd.IndexSlice[df['PD Bucket'] == 'PSI', :]
        )
        df_styled = styled_df.format(lambda x: f"{x:.4f}".rstrip('0').rstrip('.') if isinstance(x, float) else f"{x}")
        
        st.dataframe(df_styled, width=1200)
        
        data_comment_modal = Modal("Comment", key="data_comment")
        if st.button("ðŸ’¬", key="data_comment_button", help="Click here to add Comment"):
            data_comment_modal.open()
            
        # Add comment box for Data    
        if data_comment_modal.is_open():
            with data_comment_modal.container():
                data_comment_psi = st.text_area("Enter your comment:", key="data_comment_textarea")
                if st.button("Submit Comment", key="submit_data_comment"):
                    st.session_state.data_comment_psi = data_comment_psi
                    data_comment_modal.close()

        
# =============================================================================
#         # Prepare Excel file in memory
#         excel_data = io.BytesIO()
#         with pd.ExcelWriter(excel_data, engine='openpyxl') as writer:
#             df_styled.to_excel(writer, index=False, sheet_name='Calibration', startrow=12)
#         excel_data.seek(0)
#         
#         # Save byte data into st.session_state
#         st.session_state.psi_workbook_data = excel_data
#         
#         # Convert BytesIO to bytes
#         excel_data_bytes = excel_data.getvalue()
#         
#         # Encode the Excel data to base64
#         excel_data_base64 = base64.b64encode(excel_data_bytes).decode()
#         
#         # Read and encode the Excel image to base64
#         with open("excel_logo.png", "rb") as image_file:
#             excel_image_base64 = base64.b64encode(image_file.read()).decode()
#             
#         # Define the HTML for the download icon link
#         download_html = f'''
#             <a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{excel_data_base64}" 
#                download="PSI.xlsx" class="excel-download-button" title="Click here to download the Excel file">
#                <img src="data:image/png;base64,{excel_image_base64}" alt="Download Excel">
#             </a>
#         '''
#         
#         # Display the download link in Streamlit
#         st.markdown(download_html, unsafe_allow_html=True)
# =============================================================================

        # Ensure that you check if the button HTML is available in session state
        if 'excel_button_html' in st.session_state:
            st.markdown(st.session_state.excel_button_html, unsafe_allow_html=True)
        else:
            st.write("No data available for download in Excel. Please run the Overview and Change Log modules first.")
            
    with tab2:
        
        st.markdown(
                """
                <div style='text-align: center;
                font-size: 20px;'>
                    <strong>Graph</strong>
                </div>
                """,
                unsafe_allow_html=True)

        def create_subdivided_bar_plot(df):
            fig = go.Figure()
            colors = ['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd', '#8c564b', '#e377c2', '#7f7f7f', '#bcbd22', '#17becf']
            
            for i, bucket in enumerate(df['PD Bucket'].head(10).tolist()):
                values = df.iloc[i, 1:].tolist()
                years = df.columns[1:].astype(str).tolist()
                years = [i + '_' for i in years]
                
                fig.add_trace(go.Bar(
                    x=years,
                    y=values,
                    name=str(bucket),
                    marker_color=colors[i % len(colors)]
                ))
            
            fig.update_layout(
                xaxis_title='',
                yaxis_title='',
                barmode='stack',
            legend=dict(
                orientation='h', x=0.1, y=-0.2, bgcolor='rgba(255, 255, 255, 0.7)',
                bordercolor='black', borderwidth=1
            ),
            margin=dict(l=45, r=35, t=25, b=60),
            hovermode='x',
            plot_bgcolor='white',
            showlegend=True)
    
            return fig
        
        fig = create_subdivided_bar_plot(df)
        fig1_bytes = pio.to_image(fig, format='png')

        data = df.head(10)
        categories = data["PD Bucket"]
        dev_values = np.round(data["Dev"] * 100, 2)
        data.columns = [str(i) for i in data.columns]
        column_202403_values = np.round(data["202403"] * 100, 2)
        colors = ['rgb(31, 119, 180)', 'rgb(255, 127, 14)']

        trace_dev = go.Bar(
            x=[f'{cat}' for cat in categories],
            y=dev_values,
            name='Development',
            marker_color=colors[0],
            textposition='auto',
            offset=-0.2,
            width=0.4
        )
        trace_202304 = go.Bar(
            x=[f'{cat}' for cat in categories],
            y=column_202403_values,
            name='202403',
            marker_color=colors[1],
            textposition='auto',
            offset=0.2,
            width=0.4
        )
        
        fig1 = go.Figure(data=[trace_dev, trace_202304])
        fig1.update_layout(
            xaxis=dict(
                showline=True,
                linecolor='black',
                linewidth=2,
                mirror=True,
                showgrid=True,
                gridwidth=1,
                gridcolor='LightGray',
                zeroline=False,
            ),
            yaxis=dict(
                showline=True,
                linecolor='black',
                linewidth=2,
                mirror=True,
                showgrid=True,
                gridwidth=1,
                gridcolor='LightGray',
                zeroline=False,
            ),
            barmode='group',
            plot_bgcolor='white',
            legend=dict(
                orientation='h',
                x=0.4,
                y=-0.2,
                bgcolor='rgba(255, 255, 255, 0.7)',
                bordercolor='black',
                borderwidth=1,
                traceorder='normal',
            ),
            margin=dict(
                l=50,
                r=50,
                b=100,
                t=45,
                pad=4,
            ),
            height=500,
            width=1200,
        )
        fig2_bytes = pio.to_image(fig1, format='png')
        
        st.plotly_chart(fig)
        st.plotly_chart(fig1)
        
        
        # Add comment box for graph
        graph_comment_modal = Modal("Comment", key="graph_comment")
        if st.button("ðŸ’¬", key="graph_comment_button", help="Click here to add Comment"):
            graph_comment_modal.open()
            
        if graph_comment_modal.is_open():
            with graph_comment_modal.container():
                graph_comment_psi = st.text_area("Enter your comment:", key="graph_comment_textarea")
                if st.button("Submit Comment", key="submit_graph_comment"):
                    st.session_state.graph_comment_psi = graph_comment_psi
                    graph_comment_modal.close()

        st.session_state.fig1_bytes = fig1_bytes
        st.session_state.fig2_bytes = fig2_bytes
        st.session_state.thresholds_psi = thresholds_psi

        ppt_data_psi = create_powerpoint_download_button_PSI(df, fig1_bytes, fig2_bytes, thresholds_psi, st.session_state.get("data_comment_psi", ""), st.session_state.get("graph_comment_psi", ""))

        ppt_data_overview = ppt_data_change_log = ppt_data_summary = ppt_data_gini = None
        
        if 'ppt_data_overview' in st.session_state:
            ppt_data_overview = st.session_state.ppt_data_overview
        
        if 'ppt_data_change_log' in st.session_state:
            ppt_data_change_log = st.session_state.ppt_data_change_log
            
        if 'ppt_data_summary' in st.session_state:
            ppt_data_summary = st.session_state.ppt_data_summary
            
        if 'df_gini' in st.session_state and 'fig_bytes' in st.session_state and 'thresholds_gini' in st.session_state:
            df_gini = st.session_state.df_gini
            fig_bytes = st.session_state.fig_bytes
            thresholds_gini = st.session_state.thresholds_gini
            data_comment_gini = st.session_state.get("data_comment_gini", "")
            graph_comment_gini = st.session_state.get("graph_comment_gini", "")

            ppt_data_gini = create_ppt_download_button_gini(df_gini, fig_bytes, thresholds_gini, data_comment_gini, graph_comment_gini)
            
        if 'df_calibration' in st.session_state and 'fig_bytes_calibration' in st.session_state and 'thresholds_calibration' in st.session_state:
            df_calibration = st.session_state.df_calibration
            fig_bytes_calibration = st.session_state.fig_bytes_calibration
            thresholds_calibration = st.session_state.thresholds_calibration
            data_comment_calibration = st.session_state.get("data_comment_calibration", "")
            graph_comment_calibration = st.session_state.get("graph_comment_calibration", "")

            ppt_data_calibration = create_ppt_download_button_calibration(df_calibration, fig_bytes_calibration, thresholds_calibration, data_comment_calibration, graph_comment_calibration)

        if ppt_data_gini and ppt_data_calibration and ppt_data_overview and ppt_data_change_log and ppt_data_summary:
            presentation1 = load_presentation_from_bytesio(ppt_data_overview)
            presentation2 = load_presentation_from_bytesio(ppt_data_change_log)
            presentation3 = load_presentation_from_bytesio(ppt_data_summary)
            presentation4 = load_presentation_from_bytesio(ppt_data_gini)
            presentation5 = load_presentation_from_bytesio(ppt_data_calibration)
            presentation6 = load_presentation_from_bytesio(ppt_data_psi)
            
            

            merged_presentation = merge_presentations(presentation1, presentation2, presentation3, presentation4, presentation5, presentation6)
            merged_presentation_bytesio = BytesIO()
            merged_presentation.save(merged_presentation_bytesio)
            merged_presentation_bytesio.seek(0)
        
            # Base directory of the current script
            base_dir = os.path.dirname(__file__)
            
            # Construct the path for the image
            image_path = os.path.join(base_dir, "Images", "ppt_logo.png")
            
            # Read and encode the image to base64
            with open(image_path, "rb") as image_file:
                image_base64 = base64.b64encode(image_file.read()).decode()
            
            # Convert BytesIO to bytes
            ppt_data_merged_bytes = merged_presentation_bytesio.getvalue()
            
            # Encode the PPT data to base64
            ppt_data_merged_base64 = base64.b64encode(ppt_data_merged_bytes).decode()
            
            # Help text
            help_text = "Click here to download the Dashboard into PowerPoint presentation"
            
            custom_css = """
            <style>
                .ppt-download-button {
                    position: absolute;
                    top: -76px;
                    left: 85px;
                    cursor: pointer
                }
            </style>
            """
            st.markdown(custom_css, unsafe_allow_html=True)
            
            # Create the HTML for the button with image and help text
            all_ppt_button_html = f"""
            <a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{ppt_data_merged_base64}" download="Dashboard.pptx" class="ppt-download-button" title="{help_text}">
                <img src="data:image/png;base64,{image_base64}" alt="Download PPT">
            </a>
            """
            
            st.session_state.all_ppt_button_html = all_ppt_button_html
            
            # Display the custom download button in Streamlit
            st.markdown(all_ppt_button_html, unsafe_allow_html=True)

if __name__ == "__main__":
    app()